# -*- coding: utf-8 -*-
import sys

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

pptx_path = r"C:\Users\liula\Downloads\ppt2\ppt2\1.pptx"
prs = Presentation(pptx_path)


def extract_slide_data(slide, page_number):
    all_texts = []
    elements = []
    tables = []
    image_count = 0

    for shape in slide.shapes:
        try:
            shape_type = shape.shape_type
        except:
            continue

        if shape.has_text_frame:
            texts = []
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
                    all_texts.append(text)

            if texts:
                elements.append(
                    {
                        "type": "text",
                        "content": "\n".join(texts),
                        "top": int(shape.top / 914400 * 100) / 100,
                        "left": int(shape.left / 914400 * 100) / 100,
                        "width": int(shape.width / 914400 * 100) / 100,
                        "height": int(shape.height / 914400 * 100) / 100,
                    }
                )

        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            elements.append(
                {
                    "type": "image",
                    "top": int(shape.top / 914400 * 100) / 100,
                    "left": int(shape.left / 914400 * 100) / 100,
                    "width": int(shape.width / 914400 * 100) / 100,
                    "height": int(shape.height / 914400 * 100) / 100,
                }
            )

        if shape.has_table:
            table_data = []
            table = shape.table
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
            elements.append(
                {
                    "type": "table",
                    "table_rows": table_data[:6],
                    "top": int(shape.top / 914400 * 100) / 100,
                    "left": int(shape.left / 914400 * 100) / 100,
                }
            )

    return all_texts, elements, tables, image_count


def infer_layout_grid(all_texts, image_count, element_count):
    if image_count >= 2:
        return "split_2"
    elif image_count == 1 and len(all_texts) > 3:
        return "split_2"
    elif element_count > 8:
        return "grid_3"
    elif element_count > 5:
        return "grid_2"
    else:
        return "split_2"


def get_template_info(slide_type):
    templates = {
        "cover": {
            "template_family": "hero_tech_cover",
            "template_id": "hero_tech_cover",
            "skill_profile": "cover-storytelling",
            "hardness_profile": "strict",
            "schema_profile": "ppt-template/v2-hero-tech",
            "contract_profile": "cover_meta_required",
            "quality_profile": "default",
        },
        "toc": {
            "template_family": "dashboard_dark",
            "template_id": "dashboard_dark",
            "skill_profile": "dashboard-data",
            "hardness_profile": "balanced",
            "schema_profile": "ppt-template/v2-dashboard",
            "contract_profile": "chart_or_kpi_required",
            "quality_profile": "default",
        },
        "divider": {
            "template_family": "quote_hero_dark",
            "template_id": "quote_hero_dark",
            "skill_profile": "cover-default",
            "hardness_profile": "balanced",
            "schema_profile": "ppt-template/v2-quote-hero",
            "contract_profile": "cover_meta_required",
            "quality_profile": "default",
        },
        "summary": {
            "template_family": "hero_dark",
            "template_id": "hero_dark",
            "skill_profile": "cover-default",
            "hardness_profile": "balanced",
            "schema_profile": "ppt-template/v2-hero",
            "contract_profile": "default",
            "quality_profile": "default",
        },
        "content": {
            "template_family": "split_media_dark",
            "template_id": "split_media_dark",
            "skill_profile": "comparison-general",
            "hardness_profile": "balanced",
            "schema_profile": "ppt-template/v2-split-media",
            "contract_profile": "default",
            "quality_profile": "default",
        },
    }
    return templates.get(slide_type, templates["content"])


def generate_blocks(all_texts, slide_type, image_count):
    blocks = []
    card_idx = 1

    if not all_texts:
        return blocks

    blocks.append(
        {
            "block_type": "title",
            "type": "title",
            "card_id": "card-" + str(card_idx),
            "id": "card-" + str(card_idx),
            "content": all_texts[0],
            "emphasis": [],
        }
    )
    card_idx += 1

    if slide_type == "cover":
        if len(all_texts) > 1:
            blocks.append(
                {
                    "block_type": "subtitle",
                    "type": "subtitle",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": all_texts[1],
                }
            )
            card_idx += 1

        if image_count > 0:
            blocks.append(
                {
                    "block_type": "image",
                    "type": "image",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": {"url": "", "src": "", "imageUrl": ""},
                }
            )
            card_idx += 1

        for text in all_texts[2:]:
            if (
                len(text) < 100
                and "PPT" not in text
                and "your content" not in text.lower()
            ):
                blocks.append(
                    {
                        "block_type": "body",
                        "type": "body",
                        "card_id": "card-" + str(card_idx),
                        "id": "card-" + str(card_idx),
                        "content": text,
                    }
                )
                card_idx += 1
                break

    elif slide_type == "divider":
        if len(all_texts) > 1:
            blocks.append(
                {
                    "block_type": "subtitle",
                    "type": "subtitle",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": all_texts[1],
                }
            )

    elif slide_type == "toc":
        if len(all_texts) > 1:
            blocks.append(
                {
                    "block_type": "list",
                    "type": "list",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": "\n".join(all_texts[1:5]),
                }
            )

    elif slide_type == "summary":
        body_texts = [
            t
            for t in all_texts[1:]
            if len(t) < 80 and "PPT" not in t and "your content" not in t.lower()
        ]
        if body_texts:
            blocks.append(
                {
                    "block_type": "list",
                    "type": "list",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": "\n".join(body_texts[:5]),
                }
            )

    elif slide_type == "content":
        body_texts = [
            t
            for t in all_texts[1:]
            if len(t) < 80 and t and "PPT" not in t and "your content" not in t.lower()
        ]

        if body_texts:
            has_data = any("%" in t for t in body_texts)

            if has_data:
                for text in body_texts[:4]:
                    if "%" in text:
                        blocks.append(
                            {
                                "block_type": "kpi",
                                "type": "kpi",
                                "card_id": "card-" + str(card_idx),
                                "id": "card-" + str(card_idx),
                                "content": text,
                            }
                        )
                    else:
                        blocks.append(
                            {
                                "block_type": "body",
                                "type": "body",
                                "card_id": "card-" + str(card_idx),
                                "id": "card-" + str(card_idx),
                                "content": text,
                            }
                        )
                    card_idx += 1
            else:
                blocks.append(
                    {
                        "block_type": "list",
                        "type": "list",
                        "card_id": "card-" + str(card_idx),
                        "id": "card-" + str(card_idx),
                        "content": "\n".join(body_texts[:6]),
                    }
                )
                card_idx += 1

        if image_count > 0:
            blocks.append(
                {
                    "block_type": "image",
                    "type": "image",
                    "card_id": "card-" + str(card_idx),
                    "id": "card-" + str(card_idx),
                    "content": {"url": "", "src": "", "imageUrl": ""},
                }
            )

    return blocks


def generate_narration(all_texts, slide_type):
    if not all_texts:
        return ""
    title = all_texts[0]

    if slide_type == "cover":
        return (
            "欢迎各位，接下来我将为大家介绍"
            + title
            + "。本报告将从多个维度进行全面分析。"
        )
    elif slide_type == "toc":
        return "本次汇报主要分为以下几个部分，让我们逐一展开介绍。"
    elif slide_type == "divider":
        return "接下来我们进入" + title + "部分。"
    elif slide_type == "summary":
        return "以上就是本次汇报的主要内容，感谢各位的聆听。"
    else:
        return "关于" + title + "，我们来看具体内容。"


# 页码到类型的映射
page_type_hints = {
    1: "cover",
    2: "toc",
    3: "divider",
    4: "content",
    5: "content",
    6: "content",
    7: "content",
    8: "divider",
    9: "content",
    10: "content",
    11: "content",
    12: "divider",
    13: "content",
    14: "content",
    15: "content",
    16: "divider",
    17: "content",
    18: "content",
    19: "content",
    20: "summary",
}

# 提取并生成
slides_output = []
total_pages = min(20, len(prs.slides))

for idx in range(total_pages):
    slide = prs.slides[idx]
    page_num = idx + 1

    all_texts, elements, tables, image_count = extract_slide_data(slide, page_num)

    slide_type = page_type_hints.get(page_num, "content")
    if page_num == 20:
        slide_type = "summary"

    layout_grid = infer_layout_grid(all_texts, image_count, len(elements))
    template_info = get_template_info(slide_type)
    blocks = generate_blocks(all_texts, slide_type, image_count)
    narration = generate_narration(all_texts, slide_type)

    title = all_texts[0] if all_texts else ""
    body_items = [
        t
        for t in all_texts[1:]
        if len(t) < 80 and t and "PPT" not in t and "your content" not in t.lower()
    ][:6]

    slide_obj = {
        "page_number": page_num,
        "slide_id": "slide-" + str(page_num).zfill(3),
        "id": "slide-" + str(page_num).zfill(3),
        "slide_type": slide_type,
        "page_type": slide_type,
        "pageType": slide_type,
        "slideType": slide_type,
        "subtype": slide_type,
        "layout_grid": layout_grid,
        "layout": layout_grid,
        **template_info,
        "template_lock": False,
        "render_path": "svg",
        "bg_style": "light",
        "content_density": "sparse" if slide_type == "cover" else "balanced",
        "force_bento": False,
        "title": title,
        "narration": narration,
        "speaker_notes": narration,
        "speakerNotes": narration,
        "duration": 8
        if slide_type in ["cover", "summary"]
        else (5 if slide_type == "divider" else 12),
        "narration_audio_url": "",
        "narrationAudioUrl": "",
        "script": [{"role": "host", "text": narration}] if narration else [],
        "text_constraints": {
            "bullet_max_items": 6,
            "bullet_max_chars_cjk": 30,
            "min_body_font_pt": 11,
            "min_title_font_pt": 20,
            "subtitle_max_lines": 2,
            "subtitle_max_chars_cjk": 80,
            "subtitle_min_font_pt": 13,
            "bullet_auto_split": True,
        },
        "markdown": "\n".join(["# " + title] + ["- " + t for t in body_items])
        if title
        else "",
        "imageUrl": "",
        "key_points": body_items[:3],
        "bullets": body_items,
        "image_keywords": ["business", "work", "report", "summary", "corporate"],
        "blocks": blocks,
        "elements": elements[:8],
        "svg_overlay": False,
        "force_svg_overlay": False,
        "use_svg_overlay": False,
        "svg_markup": "",
    }

    slides_output.append(slide_obj)
    print(
        "已处理: 第"
        + str(page_num).zfill(2)
        + "页 ["
        + slide_type.ljust(8)
        + "] "
        + title[:25]
    )

# 构建完整输出
output = {
    "title": "202X年度工作总结PPT",
    "author": "第一PPT市场部门",
    "theme": {"palette": "business_authority", "style": "sharp"},
    "template_family": "auto",
    "template_id": "auto",
    "skill_profile": "auto",
    "hardness_profile": "balanced",
    "schema_profile": "auto",
    "contract_profile": "default",
    "quality_profile": "default",
    "design_spec": {
        "colors": {
            "primary": "C41E3A",
            "secondary": "8B0000",
            "accent": "FFFFFF",
            "background": "DC143C",
        },
        "typography": {},
        "spacing": {},
        "visual": {
            "style_recipe": "sharp",
            "visual_priority": True,
            "visual_density": "balanced",
        },
    },
    "visual_priority": True,
    "visual_preset": "executive_brief",
    "visual_density": "balanced",
    "constraint_hardness": "balanced",
    "svg_mode": "auto",
    "verbatim_content": False,
    "original_style": False,
    "disable_local_style_rewrite": False,
    "deck_id": "",
    "deck_style": "",
    "style": "",
    "minimax_palette_key": "business_authority",
    "minimax_style_variant": "sharp",
    "retry_scope": "deck",
    "retry_hint": "",
    "idempotency_key": "",
    "target_slide_ids": "",
    "target_block_ids": "",
    "slides": slides_output,
}

# 写入文件
output_path = (
    r"D:\github\with-langgraph-fastapi\test_inputs\work-summary-minimax-format.json"
)
with open(output_path, "w", encoding="utf-8") as f:
    json.dump(output, f, ensure_ascii=False, indent=2)

print("")
print("=" * 50)
print("已生成 " + str(len(slides_output)) + " 页幻灯片描述")
print("文件已保存至: " + output_path)
print("总时长: " + str(sum(s["duration"] for s in slides_output)) + " 秒")
