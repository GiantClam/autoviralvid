#!/usr/bin/env python3
"""Generate PPT v4 with yellow/gold theme to match reference PPT."""

import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Yellow/gold theme colors matching reference PPT
PRIMARY = RGBColor(0xFF, 0xC0, 0x00)  # Yellow/Gold (#FFC000)
SECONDARY = RGBColor(0x05, 0x63, 0xC1)  # Blue (#0563C1)
ACCENT = RGBColor(0xA5, 0xA5, 0xA5)  # Gray (#A5A5A5)
ORANGE = RGBColor(0xED, 0x7D, 0x31)  # Orange (#ED7D31)
DARK = RGBColor(0x44, 0x54, 0x6A)  # Dark (#44546A)
BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF)  # White
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)  # Dark gray


def create_ppt():
    """Create a simple PPT with yellow/gold theme matching reference."""
    prs = Presentation()

    # Set slide size to 13.33 x 7.5 inches (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add 20 slides
    for i in range(1, 21):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND

        # Add yellow accent bar at top
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()

        # Add blue accent bar at bottom
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = SECONDARY
        shape.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"Slide {i}"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = DARK
        run.font.name = "Microsoft YaHei"

        # Add content area with white background
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = BACKGROUND
        content_box.line.color.rgb = ACCENT
        content_box.line.width = Pt(1)

        # Add some sample content
        text_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.73), Inches(5)
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        # Title within content
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"Content for Slide {i}"
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = DARK
        run.font.name = "Microsoft YaHei"

        # Add some bullet points
        bullets = [
            "Point 1: Sample content",
            "Point 2: Business information",
            "Point 3: Key metrics",
            "Point 4: Analysis",
            "Point 5: Recommendations",
        ]

        for bullet in bullets:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"• {bullet}"
            run.font.size = Pt(18)
            run.font.color.rgb = TEXT_COLOR
            run.font.name = "Microsoft YaHei"

    # Save
    output_path = Path("output/regression/generated_v5.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Generated: {output_path}")
    return str(output_path)


if __name__ == "__main__":
    create_ppt()
