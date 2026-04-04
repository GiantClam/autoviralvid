#!/usr/bin/env python3
"""娴犲孩寮挎潻?JSON 閻㈢喐鍨?PPT閿涘牐鐨熼悽銊ョ暚閺佺繝瀵屽ù浣衡柤閿涘鈧?
閻劍纭?
    python scripts/generate_ppt_from_desc.py --input desc.json --output output.pptx
    python scripts/generate_ppt_from_desc.py --input desc.json --output output.pptx --api-url http://127.0.0.1:8124

Supported modes:
    1. API mode: call backend service (requires backend running)
    2. Local mode: call Node.js renderer directly (render fallback)
"""

from __future__ import annotations

import argparse
import base64
import shutil
import io
import json
import re
import subprocess
import sys
import time
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Dict, Optional
from urllib.parse import urlparse

import requests

REPO_ROOT = Path(__file__).resolve().parent.parent
AGENT_SRC = REPO_ROOT / "agent" / "src"
if str(AGENT_SRC) not in sys.path:
    sys.path.append(str(AGENT_SRC))

from ppt_reference_contract import audit_reference_contract

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    RGBColor = None
    Inches = None
    Pt = None
    MSO_SHAPE = None
    MSO_LINE_DASH_STYLE = None
    PP_ALIGN = None


_FOCUS_CLUSTERS = {"auto", "content", "layout", "theme", "media", "geometry"}


def _normalize_focus_cluster(value: object) -> str:
    normalized = str(value or "").strip().lower()
    return normalized if normalized in _FOCUS_CLUSTERS else "auto"


def _safe_print(message: Any) -> None:
    text = str(message or "")
    try:
        print(text)
    except UnicodeEncodeError:
        encoding = getattr(sys.stdout, "encoding", None) or "utf-8"
        fallback = text.encode(encoding, errors="replace").decode(encoding, errors="replace")
        print(fallback)


def generate_via_api(
    desc: Dict[str, Any],
    api_url: str,
    output_path: str,
    render_output_path: Optional[str] = None,
    execution_profile: str = "dev_strict",
    strict_no_fallback: bool = False,
    creation_mode: str = "fidelity",
    focus_cluster: str = "auto",
    diagnostics: Optional[Dict[str, Any]] = None,
) -> bool:
    payload = _build_pipeline_payload_from_desc(
        desc,
        execution_profile=execution_profile,
        strict_no_fallback=strict_no_fallback,
        creation_mode=creation_mode,
        focus_cluster=focus_cluster,
    )

    try:
        if render_output_path:
            req_dump_path = str(Path(render_output_path).with_suffix(".request.json"))
            Path(req_dump_path).write_text(
                json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
            )
        resp = requests.post(
            f"{api_url}/api/v1/ppt/pipeline",
            json=payload,
            headers={"Content-Type": "application/json"},
            timeout=600,
        )
        resp.raise_for_status()
        data = resp.json()

        if not data.get("success"):
            reason = f"API error: {data.get('error', 'unknown error')}"
            print(reason)
            if isinstance(diagnostics, dict):
                diagnostics["failure_stage"] = "api_pipeline"
                diagnostics["failure_reason"] = reason
            return False

        result = data.get("data", {})
        export = result.get("export", {})
        pptx_url = export.get("pptx_url") or export.get("url")
        pptx_base64 = str(export.get("pptx_base64", "") or "").strip()

        if pptx_url:
            download_url = pptx_url
            parsed = urlparse(pptx_url)
            if not parsed.scheme:
                download_url = f"{api_url}{pptx_url}"
            pptx_resp = requests.get(download_url, timeout=120)
            pptx_resp.raise_for_status()
            Path(output_path).write_bytes(pptx_resp.content)
            _safe_print(f"PPT saved to: {output_path}")

            if render_output_path:
                artifacts = result.get("artifacts", {})
                Path(render_output_path).write_text(
                    json.dumps(
                        artifacts.get("render_payload", {}),
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )
                _safe_print(f"Render payload saved to: {render_output_path}")

            return True
        elif pptx_base64:
            Path(output_path).write_bytes(base64.b64decode(pptx_base64))
            _safe_print(f"PPT saved to: {output_path}")
            if render_output_path:
                artifacts = result.get("artifacts", {})
                Path(render_output_path).write_text(
                    json.dumps(
                        artifacts.get("render_payload", {}),
                        ensure_ascii=False,
                        indent=2,
                    ),
                    encoding="utf-8",
                )
                _safe_print(f"Render payload saved to: {render_output_path}")
            return True
        else:
            reason = "API did not return PPT payload (missing pptx_url/pptx_base64)"
            print(reason)
            if isinstance(diagnostics, dict):
                diagnostics["failure_stage"] = "api_pipeline"
                diagnostics["failure_reason"] = reason
            return False

    except requests.exceptions.ConnectionError:
        reason = f"Cannot connect to API: {api_url}"
        print(reason)
        print("Please ensure backend is running (pnpm dev:agent:render)")
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "api_connectivity"
            diagnostics["failure_reason"] = reason
        return False
    except requests.exceptions.HTTPError as e:
        status_code: int | None = None
        body = ""
        if getattr(e, "response", None) is not None:
            try:
                status_code = int(e.response.status_code)
            except Exception:
                status_code = None
            body = str(e.response.text or "").strip()
        reason = f"API call failed: status={status_code or 'unknown'}"
        if body:
            reason = f"{reason}; body={body[:800]}"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "api_http_error"
            diagnostics["failure_reason"] = reason
            diagnostics["http_status"] = status_code
            diagnostics["http_body"] = body[:4000] if body else ""
        return False
    except Exception as e:
        reason = f"API call failed: {e}"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "api_exception"
            diagnostics["failure_reason"] = reason
        return False


def generate_via_local(
    desc: Dict[str, Any],
    output_path: str,
    render_output_path: Optional[str] = None,
    local_strategy: str = "reconstruct",
    reconstruct_use_template_shell: bool = False,
    reconstruct_source_aligned: bool = True,
    strict_no_fallback: bool = False,
    diagnostics: Optional[Dict[str, Any]] = None,
) -> bool:
    if local_strategy == "source-replay":
        if _generate_via_source_replay(desc, output_path, render_output_path):
            return True
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "local_source_replay"
            diagnostics["failure_reason"] = "source-replay strategy failed"
        if strict_no_fallback:
            return False
        print("source-replay strategy unavailable, fallback to reconstruct")

    if local_strategy == "reconstruct":
        if reconstruct_source_aligned:
            if _generate_via_source_aligned_reconstruct(
                desc, output_path, render_output_path
            ):
                return True
            if isinstance(diagnostics, dict):
                diagnostics["failure_stage"] = "local_reconstruct"
                diagnostics["failure_reason"] = "source-aligned reconstruct failed"
            if strict_no_fallback:
                return False
        else:
            print("reconstruct mode: source-aligned stage disabled; using pure extracted-element rebuild.")

    if _generate_from_extracted_elements(
        desc,
        output_path,
        render_output_path,
        use_template_shell=reconstruct_use_template_shell,
    ):
        return True
    if strict_no_fallback:
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "local_extracted_elements"
            diagnostics["failure_reason"] = (
                "extracted-element reconstruction failed and strict mode forbids further fallback"
            )
        return False

    scripts_dir = Path(__file__).parent
    generator_script = scripts_dir / "generate-pptx-minimax.mjs"

    if not generator_script.exists():
        reason = f"鎵句笉鍒版覆鏌撹剼鏈? {generator_script}"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "local_node_renderer"
            diagnostics["failure_reason"] = reason
        return False

    desc_path = Path(output_path).with_suffix(".desc.json")
    desc_path.write_text(
        json.dumps(desc, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    cmd = [
        "node",
        str(generator_script),
        "--input",
        str(desc_path),
        "--output",
        output_path,
    ]

    if render_output_path:
        cmd.extend(["--render-output", render_output_path])

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        if result.returncode == 0:
            _safe_print(f"PPT saved to: {output_path}")
            return True
        else:
            reason = f"娓叉煋澶辫触: {result.stderr}"
            _safe_print(reason)
            if isinstance(diagnostics, dict):
                diagnostics["failure_stage"] = "local_node_renderer"
                diagnostics["failure_reason"] = reason
            return False
    except subprocess.TimeoutExpired:
        reason = "renderer timeout"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "local_node_renderer"
            diagnostics["failure_reason"] = reason
        return False
    except Exception as e:
        reason = f"娓叉煋寮傚父: {e}"
        _safe_print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "local_node_renderer"
            diagnostics["failure_reason"] = reason
        return False


def _generate_via_source_replay(
    desc: Dict[str, Any],
    output_path: str,
    render_output_path: Optional[str] = None,
) -> bool:
    source_template = str(desc.get("source_pptx_path", "") or "").strip()
    if not source_template:
        return False
    source_path = Path(source_template)
    if not source_path.exists():
        return False

    # Current replay keeps exact fidelity by copying the full source deck.
    # It is intentionally strict: only allowed when extracted slide count
    # matches source slide count.
    try:
        if Presentation is None:
            return False
        source_prs = Presentation(str(source_path))
        source_count = len(source_prs.slides)
        selected_count = len(desc.get("slides") or [])
        if selected_count != source_count:
            return False
    except Exception:
        return False

    try:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(str(source_path), str(output_path))
        if render_output_path:
            Path(render_output_path).write_text(
                json.dumps(
                    {
                        "mode": "source_replay_copy",
                        "source_pptx_path": str(source_path.resolve()),
                        "slides": len(desc.get("slides") or []),
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
        _safe_print(f"PPT saved to: {output_path}")
        return True
    except Exception:
        return False


def _delete_slide(prs: Any, index: int) -> None:
    slide_ids = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide_id = slide_ids[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del slide_ids[index]


def _collect_selected_pages(desc: Dict[str, Any], source_count: int) -> list[int]:
    slides = desc.get("slides") if isinstance(desc.get("slides"), list) else []
    selected: list[int] = []
    seen = set()
    for idx, slide in enumerate(slides, start=1):
        if not isinstance(slide, dict):
            continue
        page_number = int(slide.get("page_number") or idx)
        if page_number < 1 or page_number > source_count:
            continue
        if page_number in seen:
            continue
        seen.add(page_number)
        selected.append(page_number)
    if selected:
        return sorted(selected)
    fallback_count = min(len(slides), source_count)
    return list(range(1, fallback_count + 1))


def _shape_anchor_inches(shape: Any) -> tuple[float, float]:
    left = float(getattr(shape, "left", 0) or 0) / 914400.0
    top = float(getattr(shape, "top", 0) or 0) / 914400.0
    return (left, top)


def _rewrite_slide_text_by_position(
    slide: Any,
    slide_desc: Dict[str, Any],
) -> int:
    max_position_distance = 0.08  # inches
    text_elements = []
    for element in slide_desc.get("elements") or []:
        if not isinstance(element, dict):
            continue
        if str(element.get("type", "")).strip().lower() != "text":
            continue
        content = str(element.get("content", "") or "").strip()
        if not content:
            continue
        text_elements.append(
            {
                "content": content,
                "left": float(element.get("left", 0) or 0),
                "top": float(element.get("top", 0) or 0),
            }
        )
    if not text_elements:
        return 0
    text_elements.sort(key=lambda item: (item["top"], item["left"]))

    text_shapes = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        left, top = _shape_anchor_inches(shape)
        text_shapes.append(
            {
                "shape": shape,
                "left": left,
                "top": top,
            }
        )
    if not text_shapes:
        return 0

    used_shapes: set[int] = set()
    changed = 0
    for element in text_elements:
        best_idx = -1
        best_dist = float("inf")
        for idx, shape_row in enumerate(text_shapes):
            if idx in used_shapes:
                continue
            dist = abs(shape_row["left"] - element["left"]) + abs(
                shape_row["top"] - element["top"]
            )
            if dist < best_dist:
                best_dist = dist
                best_idx = idx
        if best_idx < 0:
            continue
        if best_dist > max_position_distance:
            continue
        used_shapes.add(best_idx)
        target_shape = text_shapes[best_idx]["shape"]
        current = _normalize_space(str(getattr(target_shape, "text", "") or ""))
        desired = _normalize_space(element["content"])
        if current == desired:
            continue
        try:
            target_shape.text = element["content"]
            changed += 1
        except Exception:
            continue
    return changed


def _generate_via_source_aligned_reconstruct(
    desc: Dict[str, Any],
    output_path: str,
    render_output_path: Optional[str] = None,
) -> bool:
    source_template = str(desc.get("source_pptx_path", "") or "").strip()
    if not source_template or Presentation is None:
        return False
    source_path = Path(source_template)
    if not source_path.exists():
        return False

    try:
        prs = Presentation(str(source_path))
    except Exception:
        return False

    source_count = len(prs.slides)
    if source_count <= 0:
        return False

    selected_pages = _collect_selected_pages(desc, source_count)
    if not selected_pages:
        return False

    try:
        keep_indexes = {page - 1 for page in selected_pages}
        for idx in range(source_count - 1, -1, -1):
            if idx not in keep_indexes:
                _delete_slide(prs, idx)

        slide_map: Dict[int, Dict[str, Any]] = {}
        for idx, raw_slide in enumerate(desc.get("slides") or [], start=1):
            if not isinstance(raw_slide, dict):
                continue
            page_number = int(raw_slide.get("page_number") or idx)
            if page_number in selected_pages and page_number not in slide_map:
                slide_map[page_number] = raw_slide

        changed_text_count = 0
        for out_idx, page_number in enumerate(selected_pages):
            if out_idx >= len(prs.slides):
                break
            slide_desc = slide_map.get(page_number)
            if not isinstance(slide_desc, dict):
                continue
            changed_text_count += _rewrite_slide_text_by_position(
                prs.slides[out_idx], slide_desc
            )

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        if render_output_path:
            Path(render_output_path).write_text(
                json.dumps(
                    {
                        "mode": "source_aligned_reconstruct",
                        "source_pptx_path": str(source_path.resolve()),
                        "selected_pages": selected_pages,
                        "slides": len(selected_pages),
                        "changed_text_count": changed_text_count,
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
        _safe_print(f"PPT saved to: {output_path}")
        return True
    except Exception:
        return False


def _normalize_space(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _clip(text: str, max_len: int) -> str:
    value = _normalize_space(text)
    if len(value) <= max_len:
        return value
    return value[: max(0, max_len - 1)] + "..."


def _dedup_strings(items: list[str], *, limit: int) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for item in items:
        text = _normalize_space(item)
        if not text:
            continue
        key = text.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(text)
        if len(out) >= max(1, int(limit)):
            break
    return out


def _collect_reference_keywords(desc: Dict[str, Any], limit: int = 12) -> list[str]:
    tokens: list[str] = []
    stop_words = {
        "logo",
        "company",
        "slogan",
        "ppt",
        "content",
        "contents",
        "part",
        "work",
        "report",
        "http",
        "www",
    }
    for slide in desc.get("slides") or []:
        if not isinstance(slide, dict):
            continue
        title = _normalize_space(slide.get("title", ""))
        if title:
            tokens.extend(re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z][A-Za-z0-9_-]{2,}", title))
        for block in slide.get("blocks") or []:
            if not isinstance(block, dict):
                continue
            content = _normalize_space(block.get("content", ""))
            if content:
                tokens.extend(
                    re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z][A-Za-z0-9_-]{2,}", content)
                )
    dedup: list[str] = []
    seen = set()
    for raw in tokens:
        key = raw.strip().lower()
        if not key or key in seen or key in stop_words:
            continue
        seen.add(key)
        dedup.append(raw.strip())
        if len(dedup) >= limit:
            break
    return dedup


def _sanitize_desc_for_zero_create(
    desc: Dict[str, Any],
    *,
    keep_sanitized_slides: bool = False,
) -> Dict[str, Any]:
    def _norm_text(value: Any) -> str:
        if isinstance(value, str):
            return re.sub(r"\s+", " ", value).strip().lower()
        if isinstance(value, (int, float)):
            return str(value)
        if isinstance(value, dict):
            rows = []
            for key in ("title", "text", "content", "label", "description", "summary", "value"):
                text = _norm_text(value.get(key))
                if text:
                    rows.append(text)
            return " ".join(rows).strip()
        if isinstance(value, list):
            rows = [_norm_text(item) for item in value]
            return " ".join(item for item in rows if item).strip()
        return ""

    def _dedupe_key(value: str) -> str:
        return (
            str(value or "")
            .strip()
            .lower()
            .replace("\n", " ")
        )

    def _normalize_blocks(slide: Dict[str, Any]) -> list[Dict[str, Any]]:
        raw_blocks = slide.get("blocks") if isinstance(slide.get("blocks"), list) else []
        out_blocks: list[Dict[str, Any]] = []
        seen_text: set[str] = set()
        required_present = False
        required_types = {"body", "list", "quote", "icon_text", "workflow"}
        for raw in raw_blocks:
            if not isinstance(raw, dict):
                continue
            block = dict(raw)
            block_type = str(block.get("block_type") or block.get("type") or "body").strip().lower() or "body"
            block["block_type"] = block_type
            block["type"] = block_type
            content_norm = _norm_text(block.get("content")) or _norm_text(block.get("data"))
            dedupe_key = re.sub(r"\s+", " ", _dedupe_key(content_norm))
            dedupe_key = re.sub(r"[^0-9a-z\u4e00-\u9fff%+.-]", "", dedupe_key)
            if block_type != "title" and dedupe_key:
                if dedupe_key in seen_text:
                    continue
                seen_text.add(dedupe_key)
            if block_type in required_types and content_norm:
                required_present = True
            out_blocks.append(block)
        if not required_present:
            fallback_text = str(slide.get("title") or "鏍稿績瑕佺偣").strip() or "鏍稿績瑕佺偣"
            fallback_key = re.sub(r"\s+", " ", _dedupe_key(fallback_text))
            fallback_key = re.sub(r"[^0-9a-z\u4e00-\u9fff%+.-]", "", fallback_key)
            if fallback_key and fallback_key in seen_text:
                fallback_text = "鏍稿績瑕佺偣"
                fallback_key = re.sub(r"\s+", " ", _dedupe_key(fallback_text))
                fallback_key = re.sub(r"[^0-9a-z\u4e00-\u9fff%+.-]", "", fallback_key)
            if fallback_key and fallback_key in seen_text:
                page_no = int(slide.get("page_number") or 0)
                fallback_text = f"鍏抽敭瑕佺偣 {page_no or ''}".strip()
            out_blocks.append(
                {
                    "block_type": "body",
                    "type": "body",
                    "content": fallback_text,
                }
            )
        return out_blocks

    out = dict(desc if isinstance(desc, dict) else {})
    # zero_create policy:
    # - disallow direct reference-structure replay from source deck
    # - keep only prompt-level hints; let pipeline generate structured slides
    out.pop("source_pptx_path", None)
    for key in [
        "theme_manifest",
        "master_layout_manifest",
        "media_manifest",
        "theme_color_map",
    ]:
        out.pop(key, None)

    slides = out.get("slides")
    if isinstance(slides, list):
        out["requested_total_pages"] = max(3, min(50, len(slides)))
        sanitized_slides = []
        for raw in slides:
            if not isinstance(raw, dict):
                continue
            slide = dict(raw)
            for key in [
                "elements",
                "shapes",
                "media_refs",
                "slide_layout_path",
                "slide_layout_name",
                "slide_master_path",
                "slide_theme_path",
            ]:
                slide.pop(key, None)
            slide["blocks"] = _normalize_blocks(slide)
            sanitized_slides.append(slide)
        # zero_create should behave like a user prompt.
        # Keep only compact hints and drop direct slide structures.
        hints: list[str] = []
        for slide in sanitized_slides:
            title = str(slide.get("title") or "").strip()
            if title:
                hints.append(_clip(title, 90))
            for block in slide.get("blocks") or []:
                if len(hints) >= 20:
                    break
                if not isinstance(block, dict):
                    continue
                content = str(block.get("content") or "").strip()
                if content:
                    hints.append(_clip(content, 90))
            if len(hints) >= 20:
                break
        if hints:
            merged_required = out.get("required_facts")
            merged_required = merged_required if isinstance(merged_required, list) else []
            out["required_facts"] = [*merged_required, *hints][:20]
        if keep_sanitized_slides:
            out["slides"] = sanitized_slides
        else:
            out.pop("slides", None)
    return out


def _build_reconstruction_prompt(desc: Dict[str, Any]) -> str:
    slides = desc.get("slides") or []
    total_pages = len(slides)
    theme = desc.get("theme", {}) if isinstance(desc.get("theme", {}), dict) else {}
    primary = str(theme.get("primary", "") or "").strip()
    secondary = str(theme.get("secondary", "") or "").strip()
    accent = str(theme.get("accent", "") or "").strip()
    title = _clip(str(desc.get("title", "Untitled") or "Untitled"), 80)
    keywords = _collect_reference_keywords(desc, limit=10)
    keyword_text = "、".join(keywords) if keywords else "原始核心术语"
    return _clip(
        (
            f"请从零创建一份与参考风格接近的 {total_pages} 页中文商务汇报PPT。"
            f"主标题《{title}》，保持封面-目录-章节-内容-结尾叙事节奏。"
            f"配色优先 #{primary} / #{secondary} / #{accent}。"
            f"尽量保留关键词：{keyword_text}。"
        ),
        480,
    )


def _build_reference_desc_payload(desc: Dict[str, Any]) -> Dict[str, Any]:
    slides_out: list[Dict[str, Any]] = []
    for idx, slide in enumerate(desc.get("slides") or [], start=1):
        if not isinstance(slide, dict):
            continue
        slide_out: Dict[str, Any] = {
            "page_number": int(slide.get("page_number") or idx),
            "slide_id": str(slide.get("slide_id") or slide.get("id") or f"slide-{idx:03d}"),
            "id": str(slide.get("id") or slide.get("slide_id") or f"slide-{idx:03d}"),
            "slide_type": str(slide.get("slide_type") or ""),
            "title": str(slide.get("title", "") or ""),
            "layout_hint": str(slide.get("layout_hint") or ""),
            "layout_grid": str(slide.get("layout_grid") or ""),
            "render_path": str(slide.get("render_path") or ""),
            "blocks": slide.get("blocks") if isinstance(slide.get("blocks"), list) else [],
            "elements": slide.get("elements") if isinstance(slide.get("elements"), list) else [],
            "shapes": slide.get("shapes") if isinstance(slide.get("shapes"), list) else [],
            "visual": slide.get("visual") if isinstance(slide.get("visual"), dict) else {},
            "semantic_constraints": (
                slide.get("semantic_constraints")
                if isinstance(slide.get("semantic_constraints"), dict)
                else {}
            ),
            "slide_layout_path": str(slide.get("slide_layout_path") or ""),
            "slide_layout_name": str(slide.get("slide_layout_name") or ""),
            "slide_master_path": str(slide.get("slide_master_path") or ""),
            "slide_theme_path": str(slide.get("slide_theme_path") or ""),
            "media_refs": slide.get("media_refs") if isinstance(slide.get("media_refs"), list) else [],
        }
        slides_out.append(slide_out)

    out: Dict[str, Any] = {
        "title": str(desc.get("title", "") or ""),
        "source_pptx_path": str(desc.get("source_pptx_path", "") or ""),
        "theme": desc.get("theme") if isinstance(desc.get("theme"), dict) else {},
        "anchors": desc.get("anchors") if isinstance(desc.get("anchors"), list) else [],
        "required_facts": (
            desc.get("required_facts")
            if isinstance(desc.get("required_facts"), list)
            else []
        ),
        "theme_color_map": (
            desc.get("theme_color_map")
            if isinstance(desc.get("theme_color_map"), dict)
            else {}
        ),
        "media_manifest": (
            desc.get("media_manifest")
            if isinstance(desc.get("media_manifest"), list)
            else []
        ),
        "theme_manifest": (
            desc.get("theme_manifest")
            if isinstance(desc.get("theme_manifest"), list)
            else []
        ),
        "master_layout_manifest": (
            desc.get("master_layout_manifest")
            if isinstance(desc.get("master_layout_manifest"), list)
            else []
        ),
        "dimensions": (
            desc.get("dimensions") if isinstance(desc.get("dimensions"), dict) else {}
        ),
        "fonts": desc.get("fonts") if isinstance(desc.get("fonts"), list) else [],
        "slides": slides_out,
    }
    return out


def _build_pipeline_payload_from_desc(
    desc: Dict[str, Any],
    *,
    execution_profile: str = "dev_strict",
    strict_no_fallback: bool = False,
    creation_mode: str = "fidelity",
    focus_cluster: str = "auto",
) -> Dict[str, Any]:
    slides = desc.get("slides") or []
    total_pages = max(3, len(slides))
    theme = desc.get("theme", {}) if isinstance(desc.get("theme", {}), dict) else {}
    dims = desc.get("dimensions", {}) if isinstance(desc.get("dimensions", {}), dict) else {}
    width_in = float(dims.get("width_inches", 13.3333333) or 13.3333333)
    height_in = float(dims.get("height_inches", 7.5) or 7.5)
    ratio_text = f"{width_in:.2f}:{height_in:.2f}"
    normalized_creation_mode = str(creation_mode or "fidelity").strip().lower()

    topic_prompt = _build_reconstruction_prompt(desc)
    deck_title = _clip(str(desc.get("title", "Untitled") or "Untitled"), 200)
    style_pref = _clip(
        (
            f"reference-reconstruct, layout-fidelity-first, "
            f"primary=#{theme.get('primary', '')}, "
            f"secondary=#{theme.get('secondary', '')}, accent=#{theme.get('accent', '')}"
        ),
        180,
    )

    constraints: list[str] = []
    input_constraints = (
        desc.get("constraints") if isinstance(desc.get("constraints"), list) else []
    )
    for raw in input_constraints:
        text = _clip(str(raw or "").strip(), 140)
        if text:
            constraints.append(text)
    normalized_focus_cluster = _normalize_focus_cluster(focus_cluster)
    if normalized_focus_cluster != "auto":
        constraints.append(
            f"本轮优化聚焦 {normalized_focus_cluster}，其余维度保持稳定，避免跨簇大改。"
        )
    constraints.append(f"总页数固定为 {total_pages} 页，不得增减")
    constraints.append(f"画幅比例保持 {ratio_text}（16:9）")
    if theme.get("primary"):
        constraints.append(
            f"涓婚鑹插繀椤讳笌鍙傝€冩帴杩戯細primary #{theme.get('primary')}, "
            f"secondary #{theme.get('secondary')}, accent #{theme.get('accent')}"
        )
    # Add slide title anchors (up to 8) to preserve narrative skeleton.
    anchors: list[str] = []
    for idx, slide in enumerate(slides[:8], start=1):
        if not isinstance(slide, dict):
            continue
        title = _clip(str(slide.get("title", "") or "").strip(), 80)
        if title:
            constraints.append(f"第{idx}页主题锚点：{title}")
            anchors.append(title)
    anchors = _dedup_strings(anchors, limit=12)
    constraints = constraints[:20]

    required_facts: list[str] = []
    input_required_facts = (
        desc.get("required_facts")
        if isinstance(desc.get("required_facts"), list)
        else []
    )
    for raw in input_required_facts:
        text = _clip(str(raw or "").strip(), 90)
        if text:
            required_facts.append(text)
    for idx, slide in enumerate(slides, start=1):
        if len(required_facts) >= 20:
            break
        if not isinstance(slide, dict):
            continue
        title = _clip(str(slide.get("title", "") or "").strip(), 80)
        if title:
            required_facts.append(f"第{idx}页必须体现：{title}")
        for block in slide.get("blocks") or []:
            if len(required_facts) >= 20:
                break
            if not isinstance(block, dict):
                continue
            content = _clip(str(block.get("content", "") or "").strip(), 90)
            if content and len(content) >= 4:
                required_facts.append(f"保留关键文本片段：{content}")
    # de-dup
    seen_rf = set()
    rf_dedup: list[str] = []
    for item in required_facts:
        key = item.lower().strip()
        if not key or key in seen_rf:
            continue
        seen_rf.add(key)
        rf_dedup.append(item)
        if len(rf_dedup) >= 20:
            break

    palette_key = str(theme.get("palette", "auto") or "auto").strip() or "auto"
    style_variant = str(theme.get("style", "auto") or "auto").strip() or "auto"
    if style_variant not in {"auto", "sharp", "soft", "rounded", "pill"}:
        style_variant = "auto"
    template_family = (
        str(desc.get("template_family") or desc.get("template_id") or "auto").strip().lower()
        or "auto"
    )
    skill_profile = str(desc.get("skill_profile") or "auto").strip() or "auto"
    template_file_url = str(desc.get("template_file_url") or "").strip()

    if normalized_creation_mode == "zero_create":
        requested_total_pages = int(
            desc.get("requested_total_pages")
            or desc.get("total_pages")
            or len(slides)
            or 10
        )
        requested_total_pages = max(3, min(50, requested_total_pages))
        user_constraints = [
            _clip(str(item or "").strip(), 140)
            for item in (
                desc.get("constraints")
                if isinstance(desc.get("constraints"), list)
                else []
            )
            if str(item or "").strip()
        ]
        user_required_facts = [
            _clip(str(item or "").strip(), 90)
            for item in (
                desc.get("required_facts")
                if isinstance(desc.get("required_facts"), list)
                else []
            )
            if str(item or "").strip()
        ][:20]
        user_anchors = [
            _clip(str(item or "").strip(), 80)
            for item in (
                desc.get("anchors")
                if isinstance(desc.get("anchors"), list)
                else []
            )
            if str(item or "").strip()
        ]
        derived_slide_anchors: list[str] = []
        observed_block_types: set[str] = set()
        for row in slides:
            if not isinstance(row, dict):
                continue
            title = _clip(str(row.get("title") or "").strip(), 80)
            if title:
                derived_slide_anchors.append(title)
            for block in row.get("blocks") or []:
                if not isinstance(block, dict):
                    continue
                bt = str(block.get("block_type") or block.get("type") or "").strip().lower()
                if bt:
                    observed_block_types.add(bt)
        user_anchors = _dedup_strings([*user_anchors, *derived_slide_anchors], limit=12)
        template_hint_constraints: list[str] = []
        if observed_block_types & {"chart", "kpi", "table"}:
            template_hint_constraints.append("数据页优先匹配 data-capable 模板（dashboard/kpi/ops），并保留图表可读性。")
        if observed_block_types & {"workflow", "timeline", "process"}:
            template_hint_constraints.append("流程页优先匹配 process/timeline 模板，避免流程文本拥挤。")
        if observed_block_types & {"compare", "comparison", "matrix"}:
            template_hint_constraints.append("对比页优先匹配 comparison cards 模板，保持左右信息对齐。")
        if observed_block_types & {"image"}:
            template_hint_constraints.append("含图页优先匹配 split-media/image-showcase 模板，避免图文重叠。")
        template_hint_constraints.append("内容页模板需保持多样性，相邻页不得重复同一 template_family。")
        user_constraints.extend(template_hint_constraints)
        if normalized_focus_cluster != "auto":
            user_constraints.append(
                f"本轮优化目标聚焦 {normalized_focus_cluster}，其余维度保持稳定。"
            )
        user_constraints = user_constraints[:20]
        requested_quality_profile = str(desc.get("quality_profile", "auto") or "auto").strip().lower() or "auto"
        effective_quality_profile = (
            "lenient_draft" if requested_quality_profile in {"auto", "default"} else requested_quality_profile
        )
        # zero_create should favor full flow availability over strict authoring-policy failures.
        effective_execution_profile = (
            "prod_safe"
            if str(execution_profile or "").strip().lower() == "dev_strict"
            else "prod_safe"
        )
        return {
            "topic": _clip(
                str(desc.get("topic") or desc.get("title") or topic_prompt or "鍟嗕笟姹囨姤"),
                500,
            ),
            "title": deck_title,
            "audience": _clip(
                str(desc.get("audience", "general business audience") or "general business audience"),
                120,
            ),
            "purpose": _clip(
                str(
                    desc.get("purpose", "from-scratch presentation creation")
                    or "from-scratch presentation creation"
                ),
                120,
            ),
            "style_preference": _clip(
                str(desc.get("style_preference") or "professional, structured, visual-first"),
                180,
            ),
            "total_pages": requested_total_pages,
            "language": "zh-CN",
            "route_mode": "auto",
            "quality_profile": effective_quality_profile,
            "with_export": True,
            "save_artifacts": True,
            "export_channel": "local",
            # zero_create uses extracted prompt hints; skip live web enrichment
            # to avoid long-tail timeout risk in API pipeline mode.
            "web_enrichment": False,
            "image_asset_enrichment": False,
            "research_min_completeness": 0.4,
            "desired_citations": 3,
            "min_reference_materials": 3,
            "min_key_data_points": 5,
            "max_web_queries": 4,
            "max_search_results": 5,
            "minimax_style_variant": style_variant,
            "minimax_palette_key": palette_key,
            "template_family": template_family,
            "skill_profile": skill_profile,
            "template_file_url": template_file_url or None,
            "constraints": user_constraints,
            "required_facts": user_required_facts,
            "anchors": user_anchors,
            "domain_terms": _collect_reference_keywords(desc, limit=10),
            "reconstruct_from_reference": False,
            "execution_profile": effective_execution_profile,
            "force_ppt_master": False,
        }

    contract = audit_reference_contract(
        reference_desc=desc if isinstance(desc, dict) else None,
        required_facts=rf_dedup,
        anchors=anchors,
        strict=normalized_creation_mode != "zero_create",
    )
    if contract.errors:
        raise ValueError("reference contract invalid: " + "; ".join(contract.errors[:4]))

    payload = {
        "topic": topic_prompt,
        "title": deck_title,
        "audience": _clip(str(desc.get("audience", "internal business team") or "internal business team"), 120),
        "purpose": _clip(str(desc.get("purpose", "reference reconstruction") or "reference reconstruction"), 120),
        "style_preference": style_pref,
        "total_pages": total_pages,
        "language": "zh-CN",
        "route_mode": "refine",
        "quality_profile": str(desc.get("quality_profile", "auto") or "auto"),
        "with_export": True,
        "save_artifacts": True,
        "export_channel": "local",
        "web_enrichment": False,
        "research_min_completeness": 0.3,
        "desired_citations": 1,
        "min_reference_materials": 1,
        "min_key_data_points": 3,
        "max_web_queries": 1,
        "max_search_results": 3,
        "minimax_style_variant": style_variant,
        "minimax_palette_key": palette_key,
        "template_family": template_family,
        "skill_profile": skill_profile,
        "template_file_url": template_file_url or None,
        "constraints": constraints,
        "required_facts": contract.required_facts,
        "anchors": contract.anchors,
        "domain_terms": _collect_reference_keywords(desc, limit=10),
        "reconstruct_from_reference": True,
        "reference_desc": _build_reference_desc_payload(contract.reference_desc),
        "execution_profile": (
            "dev_strict"
            if str(execution_profile or "").strip().lower() == "dev_strict"
            else "prod_safe"
        ),
        "force_ppt_master": bool(strict_no_fallback or str(execution_profile or "").strip().lower() == "dev_strict"),
    }
    return payload


def _is_api_reachable(api_url: str, timeout: float = 2.0) -> bool:
    base = api_url.rstrip("/")
    probe_paths = ["/health", "/docs", "/openapi.json", "/"]
    for path in probe_paths:
        try:
            resp = requests.get(f"{base}{path}", timeout=timeout)
            if resp.status_code < 500:
                return True
        except Exception:
            continue
    return False


def _should_auto_start_local_service(api_url: str) -> bool:
    parsed = urlparse(api_url)
    host = (parsed.hostname or "").strip().lower()
    return host in {"127.0.0.1", "localhost", "0.0.0.0", "::1"}


def _start_local_render_service() -> Optional[subprocess.Popen]:
    repo_root = Path(__file__).resolve().parent.parent
    creationflags = 0
    for flag_name in ("CREATE_NEW_PROCESS_GROUP", "DETACHED_PROCESS"):
        creationflags |= int(getattr(subprocess, flag_name, 0))
    launch_candidates = [
        ["pnpm", "dev:agent:render"],
        ["pnpm.cmd", "dev:agent:render"],
        ["npm", "run", "dev:agent:render"],
    ]
    for cmd in launch_candidates:
        try:
            proc = subprocess.Popen(
                cmd,
                cwd=str(repo_root),
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                creationflags=creationflags,
            )
            print(f"Started local render service with: {' '.join(cmd)}")
            return proc
        except FileNotFoundError:
            continue
        except Exception as exc:
            print(f"Failed to launch {' '.join(cmd)}: {exc}")
            continue
    print("Failed to auto-start local render service: pnpm/npm command not found")
    return None

def _ensure_api_available(api_url: str, timeout_sec: int = 45) -> bool:
    if _is_api_reachable(api_url):
        return True

    proc = _start_local_render_service()
    if proc is None:
        return False

    deadline = time.time() + max(1, int(timeout_sec))
    while time.time() < deadline:
        if _is_api_reachable(api_url):
            return True
        if proc.poll() is not None:
            return False
        time.sleep(1.0)
    return _is_api_reachable(api_url)


def _generate_via_api_with_autostart(
    desc: Dict[str, Any],
    api_url: str,
    output_path: str,
    render_output_path: Optional[str] = None,
    execution_profile: str = "dev_strict",
    strict_no_fallback: bool = False,
    creation_mode: str = "fidelity",
    focus_cluster: str = "auto",
    diagnostics: Optional[Dict[str, Any]] = None,
) -> bool:
    success = generate_via_api(
        desc,
        api_url,
        output_path,
        render_output_path,
        execution_profile=execution_profile,
        strict_no_fallback=strict_no_fallback,
        creation_mode=creation_mode,
        focus_cluster=focus_cluster,
        diagnostics=diagnostics,
    )
    if success:
        return True
    if _is_api_reachable(api_url):
        return False
    if not _should_auto_start_local_service(api_url):
        reason = f"API is unreachable, but URL is not local; skip auto-start: {api_url}"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "api_connectivity"
            diagnostics["failure_reason"] = reason
        return False
    print("API unreachable, auto-starting local renderer: pnpm dev:agent:render")
    if not _ensure_api_available(api_url):
        reason = "Local render service failed to become healthy; cannot retry API"
        print(reason)
        if isinstance(diagnostics, dict):
            diagnostics["failure_stage"] = "api_autostart"
            diagnostics["failure_reason"] = reason
        return False
    print("Local render service is healthy; retrying API...")
    return generate_via_api(
        desc,
        api_url,
        output_path,
        render_output_path,
        execution_profile=execution_profile,
        strict_no_fallback=strict_no_fallback,
        creation_mode=creation_mode,
        focus_cluster=focus_cluster,
        diagnostics=diagnostics,
    )


def _write_failure_report(
    *,
    report_path: Path,
    execution_profile: str,
    mode: str,
    api_url: str,
    local_strategy: str,
    strict_no_fallback: bool,
    diagnostics: Dict[str, Any],
) -> None:
    report_path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "status": "failed",
        "execution_profile": execution_profile,
        "mode": mode,
        "api_url": api_url,
        "local_strategy": local_strategy,
        "strict_no_fallback": bool(strict_no_fallback),
        "failure_stage": str(diagnostics.get("failure_stage") or "unknown"),
        "failure_reason": str(diagnostics.get("failure_reason") or "unknown"),
        "timestamp": time.strftime("%Y-%m-%dT%H:%M:%S"),
    }
    report_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8"
    )


def _normalize_hex_color(raw: object) -> str:
    value = str(raw or "").strip().lstrip("#")
    if len(value) == 3 and all(c in "0123456789ABCDEFabcdef" for c in value):
        value = "".join(ch * 2 for ch in value)
    if len(value) == 6 and all(c in "0123456789ABCDEFabcdef" for c in value):
        return value.upper()
    return ""


def _resolve_color_token(raw: object, theme_color_map: Dict[str, Any]) -> str:
    token = str(raw or "").strip()
    if not token:
        return ""
    if token.lower().startswith("scheme:"):
        key = token.split(":", 1)[1].strip()
        if key:
            return _normalize_hex_color((theme_color_map or {}).get(key, ""))
        return ""
    return _normalize_hex_color(token)


def _apply_slide_background(slide: Any, bg_color_hex: str) -> None:
    if not bg_color_hex or RGBColor is None:
        return
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(bg_color_hex)
    except Exception:
        return


def _media_manifest_lookup(desc: Dict[str, Any]) -> Dict[str, bytes]:
    out: Dict[str, bytes] = {}
    manifest = desc.get("media_manifest") if isinstance(desc.get("media_manifest"), list) else []
    for row in manifest:
        if not isinstance(row, dict):
            continue
        raw_b64 = str(row.get("base64", "") or "").strip()
        path_key = str(row.get("path", "") or "").strip().replace("\\", "/")
        filename_key = str(row.get("filename", "") or "").strip()
        if not raw_b64:
            continue
        try:
            payload = base64.b64decode(raw_b64)
        except Exception:
            continue
        if not payload:
            continue
        if path_key:
            out[path_key] = payload
        if filename_key and filename_key not in out:
            out[filename_key] = payload
    return out


def _resolve_image_bytes(shape: Dict[str, Any], media_lookup: Dict[str, bytes]) -> bytes:
    image_base64 = str(shape.get("image_base64", "") or "").strip()
    if image_base64:
        try:
            payload = base64.b64decode(image_base64)
            if payload:
                return payload
        except Exception:
            pass
    media_path = str(shape.get("media_path", "") or "").strip().replace("\\", "/")
    if media_path and media_path in media_lookup:
        return media_lookup[media_path]
    media_filename = str(Path(media_path).name if media_path else "").strip()
    if media_filename and media_filename in media_lookup:
        return media_lookup[media_filename]
    return b""


def _inject_orphan_media_from_manifest(
    *,
    pptx_path: str,
    desc: Dict[str, Any],
) -> int:
    manifest = desc.get("media_manifest") if isinstance(desc.get("media_manifest"), list) else []
    if not manifest:
        return 0
    decoded: Dict[str, bytes] = {}
    for row in manifest:
        if not isinstance(row, dict):
            continue
        path = str(row.get("path", "") or "").strip().replace("\\", "/")
        raw_b64 = str(row.get("base64", "") or "").strip()
        if not path.startswith("ppt/media/") or not raw_b64:
            continue
        try:
            payload = base64.b64decode(raw_b64)
        except Exception:
            continue
        if payload:
            decoded[path] = payload
    if not decoded:
        return 0

    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            existing = {
                n
                for n in zin.namelist()
                if n.startswith("ppt/media/") and not n.endswith("/")
            }
    except Exception:
        return 0

    to_add = {path: payload for path, payload in decoded.items() if path not in existing}
    if not to_add:
        return 0
    ok = _rewrite_pptx_with_overrides(pptx_path=pptx_path, overrides=to_add)
    return len(to_add) if ok else 0


def _inject_theme_manifest(
    *,
    pptx_path: str,
    desc: Dict[str, Any],
) -> int:
    manifest = desc.get("theme_manifest") if isinstance(desc.get("theme_manifest"), list) else []
    overrides: Dict[str, bytes] = {}
    for row in manifest:
        if not isinstance(row, dict):
            continue
        path = str(row.get("path", "") or "").strip().replace("\\", "/")
        raw_b64 = str(row.get("base64", "") or "").strip()
        if not path.startswith("ppt/theme/") or not raw_b64:
            continue
        try:
            payload = base64.b64decode(raw_b64)
        except Exception:
            continue
        if payload:
            overrides[path] = payload
    if not overrides:
        return 0
    ok = _rewrite_pptx_with_overrides(pptx_path=pptx_path, overrides=overrides)
    return len(overrides) if ok else 0


def _inject_master_layout_manifest(
    *,
    pptx_path: str,
    desc: Dict[str, Any],
) -> int:
    manifest = (
        desc.get("master_layout_manifest")
        if isinstance(desc.get("master_layout_manifest"), list)
        else []
    )
    overrides: Dict[str, bytes] = {}
    for row in manifest:
        if not isinstance(row, dict):
            continue
        path = str(row.get("path", "") or "").strip().replace("\\", "/")
        raw_b64 = str(row.get("base64", "") or "").strip()
        if not raw_b64:
            continue
        if not (
            path.startswith("ppt/slideLayouts/")
            or path.startswith("ppt/slideMasters/")
        ):
            continue
        try:
            payload = base64.b64decode(raw_b64)
        except Exception:
            continue
        if payload:
            overrides[path] = payload
    if not overrides:
        return 0
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            existing = set(zin.namelist())
    except Exception:
        return 0
    to_add = {path: payload for path, payload in overrides.items() if path not in existing}
    if not to_add:
        return 0
    ok = _rewrite_pptx_with_overrides(pptx_path=pptx_path, overrides=to_add)
    return len(to_add) if ok else 0


def _rewrite_pptx_with_overrides(*, pptx_path: str, overrides: Dict[str, bytes]) -> bool:
    if not overrides:
        return True
    try:
        with tempfile.TemporaryDirectory(prefix="pptx-overrides-") as tmp_dir:
            tmp_out = Path(tmp_dir) / "repacked.pptx"
            seen: set[str] = set()
            with zipfile.ZipFile(pptx_path, "r") as zin:
                with zipfile.ZipFile(tmp_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
                    for info in zin.infolist():
                        arc_name = str(info.filename or "")
                        if not arc_name or arc_name.endswith("/"):
                            continue
                        if arc_name in overrides:
                            zout.writestr(
                                arc_name,
                                overrides[arc_name],
                                compress_type=zipfile.ZIP_DEFLATED,
                            )
                            seen.add(arc_name)
                        else:
                            zout.writestr(
                                arc_name,
                                zin.read(arc_name),
                                compress_type=info.compress_type,
                            )
                    for arc_name, payload in overrides.items():
                        if arc_name in seen:
                            continue
                        zout.writestr(
                            arc_name,
                            payload,
                            compress_type=zipfile.ZIP_DEFLATED,
                        )
            shutil.copyfile(str(tmp_out), str(pptx_path))
        return True
    except Exception:
        return False


def _generate_from_extracted_elements(
    desc: Dict[str, Any],
    output_path: str,
    render_output_path: Optional[str] = None,
    use_template_shell: bool = False,
) -> bool:
    slides = desc.get("slides", [])
    if not isinstance(slides, list) or not slides:
        return False
    if Presentation is None:
        return False

    rich_layout_signal = any(
        isinstance(s, dict)
        and isinstance(s.get("elements"), list)
        and len(s.get("elements", [])) > 0
        for s in slides
    )
    if not rich_layout_signal:
        return False

    try:
        source_template = str(desc.get("source_pptx_path", "") or "").strip()
        if use_template_shell and source_template and Path(source_template).exists():
            prs = Presentation(source_template)
            # Keep theme/master/media from template, clear slide content.
            for idx in range(len(prs.slides._sldIdLst) - 1, -1, -1):
                sld_id = prs.slides._sldIdLst[idx]
                rel_id = sld_id.rId
                prs.part.drop_rel(rel_id)
                del prs.slides._sldIdLst[idx]
        else:
            prs = Presentation()
        dims = desc.get("dimensions", {}) if isinstance(desc, dict) else {}
        width_in = float(dims.get("width_inches", 13.3333333) or 13.3333333)
        height_in = float(dims.get("height_inches", 7.5) or 7.5)
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)

        fonts = [
            str(f).strip()
            for f in (desc.get("fonts") or [])
            if str(f).strip() and not str(f).strip().startswith("+")
        ]
        if not fonts:
            fonts = ["Microsoft YaHei", "Arial"]
        theme_color_map = desc.get("theme_color_map", {}) or {}
        theme_primary_hex = _normalize_hex_color((desc.get("theme", {}) or {}).get("primary", ""))
        default_bg_hex = _normalize_hex_color((desc.get("theme", {}) or {}).get("bg", ""))
        media_lookup = _media_manifest_lookup(desc)
        align_map = {
            "l": PP_ALIGN.LEFT if PP_ALIGN is not None else None,
            "ctr": PP_ALIGN.CENTER if PP_ALIGN is not None else None,
            "r": PP_ALIGN.RIGHT if PP_ALIGN is not None else None,
            "just": PP_ALIGN.JUSTIFY if PP_ALIGN is not None else None,
            "dist": (
                getattr(PP_ALIGN, "DISTRIBUTE", None)
                if PP_ALIGN is not None
                else None
            ),
        }

        for slide_idx, slide_data in enumerate(slides):
            layout_idx = 6 if len(prs.slide_layouts) > 6 else 0
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            visual_meta = slide_data.get("visual") if isinstance(slide_data.get("visual"), dict) else {}
            bg_hex = _resolve_color_token(visual_meta.get("background_color", ""), theme_color_map)
            if not bg_hex:
                bg_hex = default_bg_hex
            _apply_slide_background(slide, bg_hex)
            text_elements = [
                el
                for el in (slide_data.get("elements") or [])
                if isinstance(el, dict) and str(el.get("type", "")).lower() == "text"
            ]
            text_elements = sorted(
                text_elements,
                key=lambda item: int(item.get("z_index", 100000) or 0),
            )

            shape_elements = slide_data.get("shapes") or []
            if isinstance(shape_elements, list):
                shape_elements = sorted(
                    [row for row in shape_elements if isinstance(row, dict)],
                    key=lambda item: int(item.get("z_index", 0) or 0),
                )
            else:
                shape_elements = []

            draw_queue: list[tuple[int, str, Dict[str, Any]]] = []
            draw_queue.extend(
                (int(shp.get("z_index", 0) or 0), "shape", shp) for shp in shape_elements
            )
            draw_queue.extend(
                (
                    int(el.get("z_index", 100000) or 0),
                    "text",
                    el,
                )
                for el in text_elements
            )
            draw_queue.sort(key=lambda row: row[0])

            for _, item_kind, item in draw_queue:
                try:
                    left = float(item.get("left", 0) or 0)
                    top = float(item.get("top", 0) or 0)
                    width = float(item.get("width", 0) or 0)
                    height = float(item.get("height", 0) or 0)
                    if width <= 0.01 or height <= 0.01:
                        continue

                    if item_kind == "shape":
                        subtype = str(item.get("subtype", "")).lower()
                        shape_kind = str(item.get("type", "")).lower()

                        image_bytes = _resolve_image_bytes(item, media_lookup)
                        if subtype == "image" or shape_kind == "image":
                            if image_bytes:
                                try:
                                    slide.shapes.add_picture(
                                        io.BytesIO(image_bytes),
                                        Inches(left),
                                        Inches(top),
                                        Inches(width),
                                        Inches(height),
                                    )
                                    continue
                                except Exception:
                                    pass

                        normalized_subtype = subtype.replace("_", "").replace("-", "")
                        if normalized_subtype in {"line", "straightconnector1"} or "connector" in normalized_subtype:
                            shape_type = getattr(MSO_SHAPE, "LINE", MSO_SHAPE.RECTANGLE)
                        elif "ellipse" in normalized_subtype or "oval" in normalized_subtype:
                            shape_type = MSO_SHAPE.OVAL
                        elif normalized_subtype in {"roundrect", "roundedrectangle"}:
                            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
                        elif normalized_subtype in {"round1rect", "round1rectangle"}:
                            shape_type = getattr(
                                MSO_SHAPE,
                                "ROUND_1_RECTANGLE",
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                            )
                        elif normalized_subtype in {"round2samerect", "round2samerectangle"}:
                            shape_type = getattr(
                                MSO_SHAPE,
                                "ROUND_2_SAME_RECTANGLE",
                                MSO_SHAPE.ROUNDED_RECTANGLE,
                            )
                        elif "parallelogram" in normalized_subtype:
                            shape_type = MSO_SHAPE.PARALLELOGRAM
                        elif normalized_subtype in {"uparrow", "uparrowcallout"}:
                            shape_type = MSO_SHAPE.UP_ARROW
                        elif "chevron" in normalized_subtype:
                            shape_type = MSO_SHAPE.CHEVRON
                        elif "teardrop" in normalized_subtype or normalized_subtype == "tear":
                            shape_type = getattr(MSO_SHAPE, "TEAR", MSO_SHAPE.OVAL)
                        elif "triangle" in normalized_subtype:
                            shape_type = getattr(
                                MSO_SHAPE,
                                "ISOSCELES_TRIANGLE",
                                MSO_SHAPE.RECTANGLE,
                            )
                        elif "donut" in normalized_subtype:
                            shape_type = getattr(MSO_SHAPE, "DONUT", MSO_SHAPE.OVAL)
                        elif normalized_subtype in {"blockarc", "arc"}:
                            shape_type = getattr(
                                MSO_SHAPE, "BLOCK_ARC", MSO_SHAPE.RECTANGLE
                            )
                        else:
                            shape_type = MSO_SHAPE.RECTANGLE
                        shape = slide.shapes.add_shape(
                            shape_type,
                            Inches(left),
                            Inches(top),
                            Inches(width),
                            Inches(height),
                        )
                        fill_hex = _resolve_color_token(item.get("fill_color"), theme_color_map)
                        line_hex = _resolve_color_token(item.get("line_color"), theme_color_map)
                        rotation = float(item.get("rotation", 0) or 0)
                        line_width_pt = float(item.get("line_width_pt", 0) or 0)
                        fill_transparency = float(item.get("fill_transparency", 0) or 0)
                        line_transparency = float(item.get("line_transparency", 0) or 0)
                        line_dash = str(item.get("line_dash", "") or "").strip().lower()

                        if fill_hex and RGBColor is not None:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
                        elif str(item.get("fill_color", "")).strip().lower() in {
                            "filled",
                            "theme",
                        } and theme_primary_hex and RGBColor is not None:
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor.from_string(theme_primary_hex)
                        else:
                            shape.fill.background()
                        if hasattr(shape.fill, "transparency"):
                            shape.fill.transparency = max(0.0, min(1.0, fill_transparency))

                        if line_hex and RGBColor is not None:
                            shape.line.fill.solid()
                            shape.line.fill.fore_color.rgb = RGBColor.from_string(line_hex)
                        else:
                            shape.line.fill.background()
                        if line_width_pt > 0:
                            shape.line.width = Pt(line_width_pt)
                        if line_dash and MSO_LINE_DASH_STYLE is not None:
                            dash_key = line_dash.replace("_", "").replace("-", "")
                            dash_map = {
                                "solid": MSO_LINE_DASH_STYLE.SOLID,
                                "dash": MSO_LINE_DASH_STYLE.DASH,
                                "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
                                "dashdotdot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
                                "longdash": MSO_LINE_DASH_STYLE.LONG_DASH,
                                "longdashdot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
                                "rounddot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                                "squaredot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
                            }
                            if dash_key in dash_map:
                                shape.line.dash_style = dash_map[dash_key]
                        if hasattr(shape.line, "transparency"):
                            shape.line.transparency = max(0.0, min(1.0, line_transparency))
                        if rotation:
                            shape.rotation = rotation
                    else:
                        content = str(item.get("content", "") or "").strip()
                        if not content:
                            continue
                        base_font_size = float(item.get("font_size", 14) or 14)
                        if base_font_size <= 0:
                            base_font_size = 14
                        text_color_hex = _resolve_color_token(
                            item.get("font_color"), theme_color_map
                        )
                        text_font_name = str(item.get("font_name", "") or "").strip()
                        text_bold = bool(item.get("bold", False))
                        text_italic = bool(item.get("italic", False))
                        align_token = str(item.get("align", "") or "").strip().lower()

                        tb = slide.shapes.add_textbox(
                            Inches(left),
                            Inches(top),
                            Inches(width),
                            Inches(height),
                        )
                        tf = tb.text_frame
                        tf.clear()
                        tf.word_wrap = True
                        tf.margin_left = 0
                        tf.margin_right = 0
                        tf.margin_top = 0
                        tf.margin_bottom = 0

                        lines = [line for line in content.split("\n")]
                        if not lines:
                            lines = [content]

                        for p_idx, line in enumerate(lines):
                            paragraph = (
                                tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                            )
                            paragraph.text = ""
                            if (
                                align_token in align_map
                                and align_map[align_token] is not None
                            ):
                                paragraph.alignment = align_map[align_token]
                            run = paragraph.add_run()
                            run.text = line
                            run.font.size = Pt(max(8, min(96, base_font_size)))
                            run.font.name = text_font_name or fonts[0]
                            run.font.bold = text_bold
                            run.font.italic = text_italic
                            if text_color_hex and RGBColor is not None:
                                run.font.color.rgb = RGBColor.from_string(text_color_hex)
                except Exception:
                    continue

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        injected_theme_count = _inject_theme_manifest(
            pptx_path=output_path,
            desc=desc,
        )
        injected_master_layout_count = _inject_master_layout_manifest(
            pptx_path=output_path,
            desc=desc,
        )
        injected_media_count = _inject_orphan_media_from_manifest(
            pptx_path=output_path,
            desc=desc,
        )

        if render_output_path:
            Path(render_output_path).write_text(
                json.dumps(
                    {
                        "mode": "extracted_element_reconstruction",
                        "slides": len(slides),
                        "font_pool_size": len(fonts),
                        "use_template_shell": bool(use_template_shell),
                        "theme_manifest_total": len(desc.get("theme_manifest") or []),
                        "theme_injected": int(injected_theme_count),
                        "master_layout_manifest_total": len(
                            desc.get("master_layout_manifest") or []
                        ),
                        "master_layout_injected": int(injected_master_layout_count),
                        "media_manifest_total": len(media_lookup),
                        "media_injected_orphan": int(injected_media_count),
                    },
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )

        _safe_print(f"PPT saved to: {output_path}")
        return True
    except Exception:
        return False


def main():
    parser = argparse.ArgumentParser(description="Generate PPT from extracted JSON.")
    parser.add_argument("--input", "-i", required=True, help="Input description JSON path")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path")
    parser.add_argument("--render-output", "-r", help="Render output JSON path")
    parser.add_argument("--api-url", default="http://127.0.0.1:8124", help="API base URL")
    parser.add_argument(
        "--mode", choices=["api", "local", "auto"], default="auto", help="Generation mode"
    )
    parser.add_argument(
        "--local-strategy",
        choices=["reconstruct", "source-replay"],
        default="reconstruct",
        help="Local generation strategy",
    )
    parser.add_argument(
        "--reconstruct-template-shell",
        choices=["on", "off"],
        default="off",
        help="Whether reconstruct strategy reuses source template shell",
    )
    parser.add_argument(
        "--reconstruct-source-aligned",
        choices=["on", "off"],
        default="on",
        help="Whether reconstruct strategy runs source-aligned replay stage first",
    )
    parser.add_argument(
        "--execution-profile",
        choices=["dev_strict", "prod_safe"],
        default="dev_strict",
        help="Execution profile: dev_strict disables fallback; prod_safe keeps fallback chain",
    )
    parser.add_argument(
        "--creation-mode",
        choices=["fidelity", "zero_create"],
        default="fidelity",
        help="fidelity may reuse reference shell in reconstruct; zero_create forbids reuse paths",
    )
    parser.add_argument(
        "--focus-cluster",
        choices=["auto", "content", "layout", "theme", "media", "geometry"],
        default="auto",
        help="Single-variable optimization focus bucket for this run",
    )
    parser.add_argument(
        "--allow-zero-create-reconstruct-overrides",
        choices=["on", "off"],
        default="off",
        help="Allow reconstruct shell/source-aligned switches in zero_create mode for controlled repair loops",
    )

    args = parser.parse_args()

    desc = json.loads(Path(args.input).read_text(encoding="utf-8"))
    creation_mode = str(args.creation_mode or "fidelity").strip().lower()
    allow_zero_create_reconstruct_overrides = (
        str(args.allow_zero_create_reconstruct_overrides or "off").strip().lower() == "on"
    )
    focus_cluster = _normalize_focus_cluster(args.focus_cluster)
    if creation_mode == "zero_create":
        # Keep sanitized slide hints for both API and local modes so zero_create
        # still carries per-page intent into skill planning and template routing.
        keep_sanitized_slides = True
        desc = _sanitize_desc_for_zero_create(
            desc,
            keep_sanitized_slides=keep_sanitized_slides,
        )
        if str(args.local_strategy) == "source-replay":
            args.local_strategy = "reconstruct"
        if not allow_zero_create_reconstruct_overrides:
            args.reconstruct_template_shell = "off"
            args.reconstruct_source_aligned = "off"
    execution_profile = str(args.execution_profile or "dev_strict").strip().lower()
    strict_no_fallback = execution_profile == "dev_strict"
    if creation_mode == "zero_create":
        # zero_create forbids direct reference-structure replay, but should still
        # allow system/model generation fallback chains.
        strict_no_fallback = False
    diagnostics: Dict[str, Any] = {}

    success = False

    if args.mode == "api":
        success = _generate_via_api_with_autostart(
            desc,
            args.api_url,
            args.output,
            args.render_output,
            execution_profile=execution_profile,
            strict_no_fallback=strict_no_fallback,
            creation_mode=creation_mode,
            focus_cluster=focus_cluster,
            diagnostics=diagnostics,
        )
    elif args.mode == "local":
        use_template_shell = str(args.reconstruct_template_shell).lower() == "on"
        use_source_aligned = str(args.reconstruct_source_aligned).lower() == "on"
        success = generate_via_local(
            desc,
            args.output,
            args.render_output,
            local_strategy=args.local_strategy,
            reconstruct_use_template_shell=use_template_shell,
            reconstruct_source_aligned=use_source_aligned,
            strict_no_fallback=strict_no_fallback,
            diagnostics=diagnostics,
        )
    else:
        use_template_shell = str(args.reconstruct_template_shell).lower() == "on"
        use_source_aligned = str(args.reconstruct_source_aligned).lower() == "on"
        success = _generate_via_api_with_autostart(
            desc,
            args.api_url,
            args.output,
            args.render_output,
            execution_profile=execution_profile,
            strict_no_fallback=strict_no_fallback,
            creation_mode=creation_mode,
            focus_cluster=focus_cluster,
            diagnostics=diagnostics,
        )
        if not success and not strict_no_fallback and creation_mode != "zero_create":
            print("API 妯″紡澶辫触锛屽皾璇曟湰鍦版ā寮?..")
            success = generate_via_local(
                desc,
                args.output,
                args.render_output,
                local_strategy=args.local_strategy,
                reconstruct_use_template_shell=use_template_shell,
                reconstruct_source_aligned=use_source_aligned,
                strict_no_fallback=False,
                diagnostics=diagnostics,
            )
        elif not success and creation_mode == "zero_create":
            print("zero_create mode: API path failed; local reconstruct fallback is disabled.")
        elif not success and strict_no_fallback:
            print("dev_strict profile: API path failed; fallback is disabled.")

    if not success:
        report_path = (
            Path(args.render_output)
            if args.render_output
            else Path(args.output).with_suffix(".failure.json")
        )
        _write_failure_report(
            report_path=report_path,
            execution_profile=execution_profile,
            mode=str(args.mode),
            api_url=str(args.api_url),
            local_strategy=str(args.local_strategy),
            strict_no_fallback=strict_no_fallback,
            diagnostics=diagnostics,
        )
        print(f"Failure report saved: {report_path}")

    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()



