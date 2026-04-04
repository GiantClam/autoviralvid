#!/usr/bin/env python3
import json
import sys
from pathlib import Path

sys.path.append(str(Path("agent/src")))


def evaluate_ppt(pptx_path):
    from ppt_quality_gate import validate_deck
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_data = {
            "page_number": i + 1,
            "slide_type": "content",
            "blocks": [],
            "title": "",
            "content": "",
        }

        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        slide_data["content"] = "\n".join(texts)
        if texts:
            slide_data["title"] = texts[0]
            slide_data["blocks"].append({"block_type": "title", "content": texts[0]})
            for text in texts[1:]:
                slide_data["blocks"].append({"block_type": "body", "content": text})

        slides.append(slide_data)

    result = validate_deck(slides, profile="default")
    return result


def main():
    generated_ppt = "output/regression/generated.pptx"
    output_dir = Path("output/regression")
    output_dir.mkdir(exist_ok=True)

    print("正在评估生成PPT的质量门控...")
    result = evaluate_ppt(generated_ppt)

    report_path = output_dir / "quality_gate_report.json"
    with open(report_path, "w", encoding="utf-8") as f:
        json.dump(result.__dict__, f, ensure_ascii=False, indent=2)

    print(f"质量门控评估完成！")
    print(f"通过: {result.passed}")
    print(f"分数: {result.score}")
    print(f"问题数量: {len(result.issues)}")
    print(f"报告已保存到: {report_path}")


if __name__ == "__main__":
    main()
