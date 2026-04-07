from pathlib import Path
import tempfile
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("test_outputs/hormuz_classroom_deck.pptx")
FONT_CN = "Microsoft YaHei"
FONT_EN = "Calibri"


COLORS = {
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
    "deep_blue": RGBColor(0, 48, 73),      # #003049
    "red": RGBColor(193, 18, 31),           # #C1121F
    "light_blue": RGBColor(102, 155, 188),  # #669BBC
    "gray_bg": RGBColor(245, 247, 251),
    "mid_gray": RGBColor(68, 84, 106),      # #44546A
}

REFERENCE_THEME = {
    "dk2": "44546A",
    "lt2": "E7E6E6",
    "accent1": "4472C4",
    "accent2": "ED7D31",
    "accent3": "A5A5A5",
    "accent4": "FFC000",
    "accent5": "5B9BD5",
    "accent6": "70AD47",
    "hlink": "0563C1",
    "folHlink": "954F72",
}


def add_full_bg(slide, color: RGBColor) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 20,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
    font_name: str | None = None,
) -> None:
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.clear()
    tf.word_wrap = bool(wrap)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    if font_name is None:
        compact = text.replace("\n", " ").strip()
        font_name = FONT_EN if compact.isascii() else FONT_CN
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.color.rgb = color or COLORS["black"]


def add_page_no(slide, page_no: int) -> None:
    add_text(
        slide,
        f"{page_no:02d}",
        x=12.2,
        y=6.8,
        w=0.8,
        h=0.3,
        size=11,
        bold=False,
        color=COLORS["mid_gray"],
        align=PP_ALIGN.RIGHT,
    )


def add_common_frame(slide, page_no: int) -> None:
    add_full_bg(slide, COLORS["white"])
    add_rect(slide, 0, 0, 13.333, 0.12, COLORS["red"])
    add_rect(slide, 0, 7.38, 13.333, 0.12, COLORS["deep_blue"])
    add_page_no(slide, page_no)


def slide_cover(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_rect(slide, 0.68, 1.65, 0.03, 3.1, COLORS["red"])
    add_rect(slide, 0.88, 1.65, 0.03, 3.1, COLORS["deep_blue"])
    add_rect(slide, 1.08, 1.65, 0.03, 3.1, COLORS["light_blue"])
    add_text(slide, "解码霍尔木兹海峡危机", 1.35, 1.85, 8.5, 0.9, size=44, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "理解其对国际关系的影响", 1.35, 2.75, 8.5, 0.7, size=28, color=COLORS["mid_gray"])
    add_rect(slide, 1.35, 3.58, 2.3, 0.03, COLORS["red"])
    add_text(slide, "大学课堂展示课件", 1.35, 3.78, 5.0, 0.5, size=20, color=COLORS["deep_blue"])
    add_text(slide, "国际关系学导论", 1.35, 4.25, 3.5, 0.4, size=16, color=COLORS["mid_gray"])
    # Decorative circles (reference-like rhythm, but not copied content)
    add_rect(slide, 10.65, 0.55, 0.9, 0.9, COLORS["light_blue"])
    add_rect(slide, 11.25, 1.05, 0.55, 0.55, COLORS["red"])
    add_rect(slide, 11.0, 4.55, 0.8, 0.8, COLORS["deep_blue"])
    add_rect(slide, 11.55, 5.05, 0.45, 0.45, COLORS["light_blue"])


def slide_contents(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_rect(slide, 0, 0, 0.16, 7.5, COLORS["deep_blue"])
    add_text(slide, "CONTENTS", 0.8, 0.45, 5.0, 0.6, size=28, bold=True, color=COLORS["deep_blue"])
    add_rect(slide, 0.8, 0.95, 2.0, 0.03, COLORS["red"])

    rows = [
        ("01", "海峡战略地位与历史背景"),
        ("02", "危机触发机制与行为体"),
        ("03", "国际影响与政策含义"),
        ("04", "总结与课堂讨论"),
    ]
    y = 1.55
    for no, title in rows:
        add_rect(slide, 0.8, y, 0.55, 0.55, COLORS["light_blue"])
        add_text(slide, no, 0.8, y + 0.07, 0.55, 0.4, size=16, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER)
        add_text(slide, title, 1.55, y + 0.04, 7.8, 0.45, size=19, color=COLORS["deep_blue"])
        y += 0.85


def slide_section(slide, page_no: int, section_no: str, part: str, title: str, subtitle: str) -> None:
    add_common_frame(slide, page_no)
    add_full_bg(slide, COLORS["gray_bg"])
    add_rect(slide, 0.7, 1.7, 4.0, 2.1, COLORS["deep_blue"])
    add_text(slide, section_no, 1.1, 2.1, 2.0, 0.7, size=62, bold=True, color=COLORS["white"])
    add_text(slide, part, 1.1, 3.0, 2.5, 0.4, size=16, bold=True, color=COLORS["light_blue"])
    add_text(slide, title, 5.3, 2.2, 7.1, 0.9, size=35, bold=True, color=COLORS["deep_blue"])
    add_text(slide, subtitle, 5.3, 3.1, 7.1, 0.6, size=21, color=COLORS["mid_gray"])


def slide_geo_value(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "海峡地理瓶颈与战略价值", 0.8, 0.45, 6.2, 0.6, size=30, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "Geographic & Strategic Importance", 0.8, 0.95, 5.8, 0.35, size=15, color=COLORS["mid_gray"])

    left_points = [
        "位于伊朗与阿曼之间，连接波斯湾与阿曼湾。",
        "通航空间有限，海上军事摩擦极易放大。",
        "国际能源运输高度依赖该海上通道。",
        "区域危机可迅速外溢为全球市场冲击。",
    ]
    y = 1.85
    for row in left_points:
        add_rect(slide, 0.95, y, 0.08, 0.22, COLORS["red"])
        add_text(slide, row, 1.15, y - 0.02, 5.1, 0.35, size=14, color=COLORS["deep_blue"])
        y += 0.95

    add_text(slide, "关键观察", 7.1, 1.8, 3.8, 0.45, size=18, bold=True, color=COLORS["deep_blue"])
    add_rect(slide, 7.1, 2.35, 4.7, 1.1, COLORS["gray_bg"])
    add_text(slide, "约 1,700 万桶/日", 7.35, 2.62, 3.0, 0.4, size=24, bold=True, color=COLORS["red"])
    add_text(slide, "经该航道运输（课堂口径，非实时交易数据）", 7.35, 3.0, 4.2, 0.3, size=11, color=COLORS["mid_gray"])
    add_rect(slide, 7.1, 3.75, 4.7, 1.1, COLORS["gray_bg"])
    add_text(slide, "约 20% 海运原油贸易关联", 7.35, 4.02, 4.2, 0.4, size=20, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "风险主要体现在价格、保险与航线冗余成本", 7.35, 4.37, 4.2, 0.3, size=11, color=COLORS["mid_gray"])


def slide_timeline(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "紧张局势演变：关键节点回顾", 0.7, 0.42, 8.6, 0.6, size=28, bold=True, color=COLORS["deep_blue"])
    rows = [
        ("1979", "革命后美伊关系恶化，地区安全结构重排。"),
        ("1980-88", "两伊战争期间，海上通道安全多次受到冲击。"),
        ("2012", "制裁升级，围绕海峡通行权的政治威慑增强。"),
        ("2019", "油轮与无人机事件叠加，误判风险上升。"),
        ("2023+", "危机高强度波动常态化，降温机制仍不稳定。"),
    ]
    y = 1.35
    for year, text in rows:
        add_rect(slide, 1.2, y + 0.08, 0.14, 0.14, COLORS["red"])
        add_text(slide, year, 1.6, y, 1.6, 0.35, size=15, bold=True, color=COLORS["deep_blue"])
        add_text(slide, text, 3.0, y, 8.8, 0.35, size=13, color=COLORS["mid_gray"])
        add_rect(slide, 1.27, y + 0.22, 0.01, 0.58, COLORS["light_blue"])
        y += 1.0


def slide_stakeholders(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "主要行为体与策略偏好", 0.65, 0.35, 7.0, 0.6, size=28, bold=True, color=COLORS["deep_blue"])
    headers = ["伊朗", "美国及伙伴", "海湾国家"]
    bodies = [
        "强调主权安全与反制筹码，倾向不对称威慑。",
        "强调航道自由与盟友保护，提升军事存在强度。",
        "在安全依赖与经济稳定之间寻找政策平衡。",
    ]
    xs = [0.7, 4.25, 7.8]
    for idx, x in enumerate(xs):
        add_rect(slide, x, 1.55, 3.0, 0.45, COLORS["deep_blue"])
        add_text(slide, headers[idx], x + 0.1, 1.62, 2.8, 0.3, size=15, bold=True, color=COLORS["white"])
        add_rect(slide, x, 2.0, 3.0, 4.1, COLORS["gray_bg"])
        add_text(slide, bodies[idx], x + 0.18, 2.22, 2.65, 1.45, size=12, color=COLORS["deep_blue"])
    add_text(slide, "共同约束：都希望避免全面失控，但对“可接受风险”认知差异显著。", 0.85, 6.15, 11.8, 0.45, size=14, color=COLORS["mid_gray"])


def slide_causes(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "危机根源：五个相互强化的驱动因素", 0.65, 0.35, 9.8, 0.6, size=26, bold=True, color=COLORS["deep_blue"])
    rows = [
        ("核问题争议", "核计划与制裁循环使互信持续下降。"),
        ("地区权力竞争", "伊朗与海湾国家围绕地区影响力博弈。"),
        ("身份与叙事", "宗教与历史记忆加剧安全焦虑。"),
        ("外部力量介入", "域外军事存在改变地区威慑平衡。"),
        ("能源与金融联动", "油价、保险和航运成本共同放大冲击。"),
    ]
    y = 1.4
    for i, (k, v) in enumerate(rows):
        color = COLORS["light_blue"] if i % 2 == 0 else RGBColor(228, 236, 246)
        add_rect(slide, 0.8, y, 11.8, 0.88, color)
        add_text(slide, k, 1.0, y + 0.22, 2.2, 0.3, size=14, bold=True, color=COLORS["deep_blue"])
        add_text(slide, v, 3.1, y + 0.22, 9.2, 0.3, size=13, color=COLORS["mid_gray"])
        y += 1.02


def slide_energy_impact(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "对全球能源安全的冲击路径", 0.8, 0.35, 5.8, 0.6, size=28, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "Impact on Global Energy Security", 0.8, 0.82, 5.2, 0.3, size=14, color=COLORS["mid_gray"])

    left = [
        ("油价波动", "地缘风险溢价抬升，市场预期先于实物流动反应。"),
        ("供应风险", "关键通道受阻将导致运输时滞与库存压力。"),
        ("政策反应", "进口国加大储备、长期合同与运输分流布局。"),
    ]
    y = 1.55
    for title, detail in left:
        add_text(slide, title, 0.8, y, 2.4, 0.35, size=15, bold=True, color=COLORS["deep_blue"])
        add_text(slide, detail, 0.8, y + 0.28, 5.7, 0.32, size=12, color=COLORS["mid_gray"])
        y += 1.55

    add_rect(slide, 6.2, 1.55, 6.1, 1.35, COLORS["gray_bg"])
    add_text(slide, "替代路线建设", 6.45, 1.8, 2.2, 0.35, size=15, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "陆上管道与港口配套可降低单点依赖。", 6.45, 2.12, 5.5, 0.3, size=12, color=COLORS["mid_gray"])

    add_rect(slide, 6.2, 3.15, 6.1, 1.35, COLORS["gray_bg"])
    add_text(slide, "战略储备机制", 6.45, 3.4, 2.2, 0.35, size=15, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "短期缓冲有效，但无法替代长期地缘风险治理。", 6.45, 3.72, 5.5, 0.3, size=12, color=COLORS["mid_gray"])

    add_text(slide, "结论：能源安全问题正从“市场议题”转为“安全议题”。", 0.8, 6.25, 11.4, 0.35, size=13, bold=True, color=COLORS["red"])


def slide_ir_impact(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "对国际关系格局的系统影响", 0.65, 0.28, 7.5, 0.6, size=27, bold=True, color=COLORS["deep_blue"])
    cards = [
        ("美伊关系", "“施压-反制-降温”循环反复，长期互信难以重建。"),
        ("地区阵营化", "安全联盟与代理冲突并存，区域分化加深。"),
        ("大国博弈", "域外大国在中东的能源与安全利益交叉竞争。"),
        ("规则挑战", "航行自由与主权安全叙事在制度层面持续冲突。"),
    ]
    positions = [(0.8, 1.35), (4.4, 1.35), (8.0, 1.35), (0.8, 4.0)]
    widths = [3.3, 3.3, 4.3, 7.5]
    for idx, (title, body) in enumerate(cards):
        x, y = positions[idx]
        w = widths[idx]
        add_rect(slide, x, y, w, 0.42, COLORS["deep_blue"])
        add_text(slide, title, x + 0.12, y + 0.07, w - 0.2, 0.28, size=14, bold=True, color=COLORS["white"])
        add_rect(slide, x, y + 0.42, w, 1.5 if idx < 3 else 1.2, COLORS["gray_bg"])
        add_text(slide, body, x + 0.14, y + 0.66, w - 0.3, 0.9, size=12, color=COLORS["mid_gray"])
    add_rect(slide, 8.0, 4.0, 4.3, 1.62, COLORS["gray_bg"])
    add_text(slide, "课堂延伸：若缺乏危机沟通机制，偶发事件更容易升级为制度性对抗。", 8.2, 4.3, 3.9, 1.0, size=12, color=COLORS["red"])


def slide_summary(slide, page_no: int) -> None:
    add_common_frame(slide, page_no)
    add_text(slide, "总结与思考", 0.65, 0.25, 4.0, 0.6, size=29, bold=True, color=COLORS["deep_blue"])
    add_text(slide, "核心要点", 0.8, 1.2, 3.0, 0.35, size=17, bold=True, color=COLORS["deep_blue"])
    add_rect(slide, 0.8, 1.55, 5.5, 2.2, COLORS["gray_bg"])
    add_text(
        slide,
        "霍尔木兹海峡危机本质是地缘安全、能源通道与大国政治交织。\n稳定不只依赖威慑，更依赖透明沟通与可验证规则。",
        1.05,
        1.85,
        5.0,
        1.5,
        size=13,
        color=COLORS["mid_gray"],
    )

    add_text(slide, "讨论问题", 6.8, 1.2, 3.0, 0.35, size=17, bold=True, color=COLORS["deep_blue"])
    add_rect(slide, 6.8, 1.55, 5.7, 2.2, COLORS["gray_bg"])
    add_text(
        slide,
        "1) 如何降低“防御行为被误判”为进攻准备？\n2) 能源转型能否显著缓解该海峡的战略敏感性？",
        7.05,
        1.85,
        5.2,
        1.5,
        size=13,
        color=COLORS["mid_gray"],
    )
    add_rect(slide, 0, 4.7, 13.333, 1.1, COLORS["deep_blue"])
    add_text(
        slide,
        "感谢聆听",
        0,
        4.85,
        13.333,
        0.34,
        size=21,
        bold=True,
        color=COLORS["white"],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Thank You",
        0,
        5.33,
        13.333,
        0.24,
        size=13,
        color=COLORS["light_blue"],
        align=PP_ALIGN.CENTER,
        font_name=FONT_EN,
    )


def patch_theme_to_reference(pptx_path: Path) -> None:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    with zipfile.ZipFile(pptx_path, "r") as zin:
        infos = zin.infolist()
        content = {info.filename: zin.read(info.filename) for info in infos}

    theme_key = "ppt/theme/theme1.xml"
    theme_bytes = content.get(theme_key)
    if theme_bytes:
        try:
            root = ET.fromstring(theme_bytes)
            clr_scheme = root.find(".//a:clrScheme", ns)
            if clr_scheme is not None:
                for node in list(clr_scheme):
                    tag = node.tag.split("}")[-1]
                    if tag in REFERENCE_THEME:
                        node.clear()
                        srgb = ET.SubElement(
                            node,
                            "{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr",
                        )
                        srgb.set("val", REFERENCE_THEME[tag])

            major_latin = root.find(".//a:majorFont/a:latin", ns)
            minor_latin = root.find(".//a:minorFont/a:latin", ns)
            if major_latin is not None:
                major_latin.set("typeface", "Calibri Light")
            if minor_latin is not None:
                minor_latin.set("typeface", "Calibri")
            content[theme_key] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        except Exception:
            pass

    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=".pptx",
        dir=str(pptx_path.parent),
    ) as tmp:
        tmp_path = Path(tmp.name)
    try:
        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in infos:
                data = content.get(info.filename, b"")
                zout.writestr(info, data)
        tmp_path.replace(pptx_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink(missing_ok=True)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs.slides.add_slide(prs.slide_layouts[6]), 1)
    slide_contents(prs.slides.add_slide(prs.slide_layouts[6]), 2)
    slide_section(prs.slides.add_slide(prs.slide_layouts[6]), 3, "01", "PART 1", "海峡战略地位", "地理位置与历史背景")
    slide_geo_value(prs.slides.add_slide(prs.slide_layouts[6]), 4)
    slide_timeline(prs.slides.add_slide(prs.slide_layouts[6]), 5)
    slide_section(prs.slides.add_slide(prs.slide_layouts[6]), 6, "02", "PART 2", "危机触发机制", "行为体动机与风险放大")
    slide_stakeholders(prs.slides.add_slide(prs.slide_layouts[6]), 7)
    slide_causes(prs.slides.add_slide(prs.slide_layouts[6]), 8)
    slide_section(prs.slides.add_slide(prs.slide_layouts[6]), 9, "03", "PART 3", "国际影响", "能源安全与秩序竞争")
    slide_energy_impact(prs.slides.add_slide(prs.slide_layouts[6]), 10)
    slide_ir_impact(prs.slides.add_slide(prs.slide_layouts[6]), 11)
    slide_summary(prs.slides.add_slide(prs.slide_layouts[6]), 12)
    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(str(OUT_PATH))
    patch_theme_to_reference(OUT_PATH)
    print(str(OUT_PATH.resolve()))


if __name__ == "__main__":
    main()
