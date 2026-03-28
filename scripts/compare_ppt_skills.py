"""
Compare MiniMax PPTX style variants on the same input payload.

Usage:
  python scripts/compare_ppt_skills.py --input <export_req.json> [--output-dir <dir>]
"""

from __future__ import annotations

import argparse
import json
import re
import subprocess
import time
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
MINIMAX_GEN = ROOT / "scripts" / "generate-pptx-minimax.mjs"


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def strip_html(text: str) -> str:
    out = re.sub(r"<script[\s\S]*?</script>", " ", text or "", flags=re.I)
    out = re.sub(r"<[^>]+>", " ", out)
    out = out.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    return normalize_text(out)


def read_payload(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if "slides" not in payload or not isinstance(payload["slides"], list):
        raise ValueError("input JSON must contain 'slides' array")
    return payload


def expected_text_stats(payload: dict[str, Any]) -> tuple[int, int]:
    slides = payload.get("slides", [])
    expected_titles = 0
    expected_chars = 0
    for slide in slides:
        title = normalize_text(str(slide.get("title", "")))
        if title:
            expected_titles += 1
            expected_chars += len(title)
        for el in slide.get("elements", []):
            if str(el.get("type", "")).lower() != "text":
                continue
            expected_chars += len(strip_html(str(el.get("content", ""))))
        expected_chars += len(normalize_text(str(slide.get("narration", ""))))
    return expected_titles, expected_chars


def extract_ppt_stats(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            val = normalize_text(shape.text_frame.text or "")
            if val:
                texts.append(val)
    joined = "\n".join(texts)
    return {
        "slide_count": len(prs.slides),
        "text_chars": len(joined),
        "all_text": joined,
    }


def score_variant(
    name: str,
    expected_slide_count: int,
    expected_titles: int,
    expected_chars: int,
    output_path: Path,
    elapsed_ms: int,
    source_titles: list[str],
) -> dict[str, Any]:
    stats = extract_ppt_stats(output_path)
    all_text = stats["all_text"]
    title_hits = 0
    for title in source_titles:
        if title and title in all_text:
            title_hits += 1

    slide_score = 40 if stats["slide_count"] == expected_slide_count else max(
        0, 40 - 10 * abs(stats["slide_count"] - expected_slide_count)
    )
    coverage_ratio = min(1.0, stats["text_chars"] / max(1, expected_chars))
    coverage_score = 35 * coverage_ratio
    title_ratio = title_hits / max(1, expected_titles)
    title_score = 25 * title_ratio
    total_score = round(slide_score + coverage_score + title_score, 2)

    return {
        "name": name,
        "output": str(output_path),
        "elapsed_ms": elapsed_ms,
        "size_bytes": output_path.stat().st_size if output_path.exists() else 0,
        "slide_count": stats["slide_count"],
        "text_chars": stats["text_chars"],
        "title_hits": title_hits,
        "score": total_score,
        "score_breakdown": {
            "slide_score": round(slide_score, 2),
            "coverage_score": round(coverage_score, 2),
            "title_score": round(title_score, 2),
            "coverage_ratio": round(coverage_ratio, 4),
            "title_ratio": round(title_ratio, 4),
        },
    }


def run_generator(cmd: list[str], cwd: Path, timeout_sec: int) -> tuple[int, str]:
    started = time.perf_counter()
    proc = subprocess.run(
        cmd,
        cwd=str(cwd),
        capture_output=True,
        text=True,
        timeout=timeout_sec,
    )
    elapsed_ms = int((time.perf_counter() - started) * 1000)
    if proc.returncode != 0:
        raise RuntimeError(f"command failed: {' '.join(cmd)}\n{proc.stderr or proc.stdout}")
    return elapsed_ms, proc.stdout.strip()


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True, help="Path to export_req.json")
    parser.add_argument(
        "--output-dir",
        default="",
        help="Directory for generated artifacts",
    )
    parser.add_argument("--timeout", type=int, default=120)
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"input not found: {input_path}")

    if args.output_dir:
        output_dir = Path(args.output_dir).resolve()
    else:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_dir = (ROOT / "test_outputs" / f"ppt_skill_compare_{stamp}").resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    payload = read_payload(input_path)
    expected_slide_count = len(payload["slides"])
    expected_titles, expected_chars = expected_text_stats(payload)
    source_titles = [normalize_text(str(s.get("title", ""))) for s in payload.get("slides", [])]

    minimax_auto_out = output_dir / "minimax_auto.pptx"
    minimax_sharp_out = output_dir / "minimax_sharp.pptx"

    minimax_auto_ms, _ = run_generator(
        ["node", str(MINIMAX_GEN), "--input", str(input_path), "--output", str(minimax_auto_out), "--style", "auto"],
        ROOT,
        args.timeout,
    )
    minimax_sharp_ms, _ = run_generator(
        ["node", str(MINIMAX_GEN), "--input", str(input_path), "--output", str(minimax_sharp_out), "--style", "sharp"],
        ROOT,
        args.timeout,
    )

    minimax_auto_result = score_variant(
        "minimax_auto",
        expected_slide_count,
        expected_titles,
        expected_chars,
        minimax_auto_out,
        minimax_auto_ms,
        source_titles,
    )
    minimax_sharp_result = score_variant(
        "minimax_sharp",
        expected_slide_count,
        expected_titles,
        expected_chars,
        minimax_sharp_out,
        minimax_sharp_ms,
        source_titles,
    )

    winner = (
        minimax_auto_result
        if minimax_auto_result["score"] >= minimax_sharp_result["score"]
        else minimax_sharp_result
    )
    report = {
        "input": str(input_path),
        "expected_slide_count": expected_slide_count,
        "expected_titles": expected_titles,
        "expected_chars": expected_chars,
        "results": [minimax_auto_result, minimax_sharp_result],
        "winner": winner["name"],
        "generated_at": datetime.now().isoformat(),
    }

    report_path = output_dir / "comparison.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    print(json.dumps(report, ensure_ascii=False, indent=2))
    print(f"\nSaved comparison report: {report_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
