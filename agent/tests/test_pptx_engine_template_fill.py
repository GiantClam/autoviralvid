import io
import zipfile
import base64

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src import pptx_engine
from src.pptx_engine import fill_template_pptx


def _mock_markitdown(monkeypatch, text: str):
    monkeypatch.setenv("PPT_TEMPLATE_MARKITDOWN_ENABLED", "true")

    def _fake_extract_text_with_markitdown(_pptx_bytes, timeout_sec=20):
        assert timeout_sec >= 5
        return {
            "enabled": True,
            "ok": True,
            "error": "",
            "text": text,
            "text_length": len(text),
        }

    monkeypatch.setattr("src.ppt_visual_qa.extract_text_with_markitdown", _fake_extract_text_with_markitdown)


def _make_template_with_tokens() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, prs.slide_width, 600000).text_frame.text = "{{deck_title}}"
    slide.shapes.add_textbox(0, 700000, prs.slide_width, 600000).text_frame.text = "{{slide_1_title}}"
    slide.shapes.add_textbox(0, 1400000, prs.slide_width, 600000).text_frame.text = "{{body}}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_fill_template_pptx_replaces_placeholders(monkeypatch):
    template_bytes = _make_template_with_tokens()
    _mock_markitdown(monkeypatch, "## Intro Slide\n- Point A\n- Point B")
    result = fill_template_pptx(
        template_bytes=template_bytes,
        slides=[
            {
                "slide_id": "s1",
                "title": "Intro Slide",
                "elements": [
                    {"type": "text", "content": "Point A"},
                    {"type": "text", "content": "Point B"},
                ],
                "narration": "Narration line",
            }
        ],
        deck_title="Demo Deck",
        author="AutoViralVid",
    )
    assert isinstance(result.get("pptx_bytes"), (bytes, bytearray))
    assert int(result.get("replacement_count") or 0) >= 3
    prs = Presentation(io.BytesIO(result["pptx_bytes"]))
    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    merged = "\n".join(texts)
    assert "Demo Deck" in merged
    assert "Intro Slide" in merged
    assert "Point A" in merged


def test_fill_template_pptx_defaults_to_xml_edit_engine(monkeypatch):
    monkeypatch.delenv("PPT_TEMPLATE_EDIT_ENGINE", raising=False)
    template_bytes = _make_template_with_tokens()
    _mock_markitdown(monkeypatch, "## Intro Slide\n- Point A")
    result = fill_template_pptx(
        template_bytes=template_bytes,
        slides=[
            {
                "slide_id": "s1",
                "title": "Intro Slide",
                "elements": [{"type": "text", "content": "Point A"}],
                "narration": "Narration line",
            }
        ],
        deck_title="Demo Deck",
        author="AutoViralVid",
    )
    assert result.get("engine") == "xml"
    assert int(result.get("replacement_count") or 0) >= 3


def _make_template_with_image_token() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    box.text_frame.text = "{{image}}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_template_with_orphan_media() -> bytes:
    base = _make_template_with_tokens()
    src = io.BytesIO(base)
    out = io.BytesIO()
    orphan_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAukB9m6lJgAAAABJRU5ErkJggg=="
    )
    with zipfile.ZipFile(src, mode="r") as zin:
        with zipfile.ZipFile(out, mode="w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                zout.writestr(info, zin.read(info.filename))
            zout.writestr("ppt/media/orphan_unused.png", orphan_png)
    return out.getvalue()


def test_fill_template_pptx_replaces_image_placeholder(monkeypatch):
    template_bytes = _make_template_with_image_token()
    one_px_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAukB9m6lJgAAAABJRU5ErkJggg=="
    )
    monkeypatch.setattr(pptx_engine, "_download_image", lambda _url: one_px_png)
    monkeypatch.setenv("PPT_TEMPLATE_EDIT_ENGINE", "xml")
    _mock_markitdown(monkeypatch, "## Intro\n- Visual")

    result = fill_template_pptx(
        template_bytes=template_bytes,
        slides=[{"slide_id": "s1", "title": "Intro", "image_url": "https://example.com/demo.png"}],
        deck_title="Demo Deck",
        author="AutoViralVid",
    )

    assert result.get("engine") == "python-pptx"
    assert int(result.get("image_replacement_count") or 0) >= 1
    prs = Presentation(io.BytesIO(result["pptx_bytes"]))
    shape_types = [shape.shape_type for shape in prs.slides[0].shapes]
    assert MSO_SHAPE_TYPE.PICTURE in shape_types


def test_fill_template_xml_cleans_orphan_media(monkeypatch):
    monkeypatch.setenv("PPT_TEMPLATE_EDIT_ENGINE", "xml")
    template_bytes = _make_template_with_orphan_media()
    _mock_markitdown(monkeypatch, "## Intro Slide\n- Point A")

    result = fill_template_pptx(
        template_bytes=template_bytes,
        slides=[{"slide_id": "s1", "title": "Intro Slide", "elements": [{"type": "text", "content": "Point A"}]}],
        deck_title="Demo Deck",
        author="AutoViralVid",
    )

    assert result.get("engine") == "xml"
    assert int(result.get("cleaned_resource_count") or 0) >= 1
    with zipfile.ZipFile(io.BytesIO(result["pptx_bytes"]), mode="r") as zf:
        names = set(zf.namelist())
    assert "ppt/media/orphan_unused.png" not in names


def test_fill_template_uses_markitdown_structure_for_replacements(monkeypatch):
    template_bytes = _make_template_with_tokens()
    monkeypatch.setenv("PPT_TEMPLATE_MARKITDOWN_ENABLED", "true")

    def _fake_extract_text_with_markitdown(_pptx_bytes, timeout_sec=20):
        assert timeout_sec >= 5
        return {
            "enabled": True,
            "ok": True,
            "error": "",
            "text": "## Template Parsed Title\n- Insight A\n- Insight B",
            "text_length": 48,
        }

    monkeypatch.setattr("src.ppt_visual_qa.extract_text_with_markitdown", _fake_extract_text_with_markitdown)

    result = fill_template_pptx(
        template_bytes=template_bytes,
        slides=[{"slide_id": "s1", "title": "", "elements": [], "narration": ""}],
        deck_title="Demo Deck",
        author="AutoViralVid",
    )

    assert result.get("markitdown_enabled") is True
    assert result.get("markitdown_ok") is True
    assert result.get("markitdown_used") is True
    prs = Presentation(io.BytesIO(result["pptx_bytes"]))
    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    merged = "\n".join(texts)
    assert "Template Parsed Title" in merged
    assert "Insight A" in merged
