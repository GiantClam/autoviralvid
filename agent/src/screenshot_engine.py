"""
PPT 截图引擎 v2 — 6种完全不同模板 + 充实内容 + 高亮重点

修复:
1. 每页模板完全不同 (封面/要点/对比/数据/金句/致谢)
2. 内容充实，减少留白
3. 高亮重点 (红色标注+大数字)
"""

from __future__ import annotations
import json, logging, os, subprocess
from typing import List

logger = logging.getLogger("screenshot_v2")

W, H = 1920, 1080


def _highlight(text: str, emphasis: list) -> str:
    """高亮重点词汇"""
    for w in emphasis:
        if w and w in text:
            text = text.replace(
                w,
                f'<span style="color:#ef4444;font-weight:900;text-decoration:underline;text-decoration-color:#ef4444;text-underline-offset:4px;">{w}</span>',
            )
    return text


# ════════════════════════════════════════════════════════════════════
# 6种完全不同的模板
# ════════════════════════════════════════════════════════════════════


def _render_cover(slide: dict) -> str:
    """模板1: 封面 — 全屏渐变+超大标题"""
    c = slide.get("content", {})
    title = c.get("title", "")
    items = c.get("body_items", [])
    subtitle = c.get("subtitle", "") or (items[0] if items else "")
    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e3a5f 40%,#0f172a 100%);color:#e2e8f0;font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;position:relative}}
.circle{{position:absolute;width:600px;height:600px;border-radius:50%;background:radial-gradient(circle,rgba(56,189,248,0.08),transparent 70%);top:-100px;right:-100px}}
.circle2{{position:absolute;width:400px;height:400px;border-radius:50%;background:radial-gradient(circle,rgba(139,92,246,0.06),transparent 70%);bottom:-80px;left:-80px}}
h1{{font-size:72px;font-weight:900;background:linear-gradient(135deg,#38bdf8,#a78bfa);-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:20px;line-height:1.1}}
.line{{width:80px;height:4px;background:linear-gradient(90deg,#38bdf8,#8b5cf6);border-radius:2px;margin:0 auto 20px}}
p{{font-size:26px;color:#94a3b8;line-height:1.6;max-width:800px}}
</style></head><body>
<div class="circle"></div><div class="circle2"></div>
<h1>{title}</h1><div class="line"></div><p>{subtitle}</p>
</body></html>"""


def _render_bullet_points(slide: dict) -> str:
    """模板2: 要点列表 — 顶部标题栏+要点+右侧高亮"""
    c = slide.get("content", {})
    title = c.get("title", "")
    items = c.get("body_items", [])
    emphasis = c.get("emphasis_words", [])

    items_html = ""
    for item in items:
        highlighted = _highlight(item, emphasis)
        items_html += f'<div style="display:flex;gap:16px;padding:14px 20px;background:rgba(255,255,255,0.04);border-radius:10px;border-left:4px solid #38bdf8;margin-bottom:10px;"><div style="width:8px;height:8px;border-radius:50%;background:#38bdf8;margin-top:8px;flex-shrink:0;"></div><div style="font-size:24px;line-height:1.5;color:#e2e8f0;">{highlighted}</div></div>'

    emphasis_box = ""
    if emphasis:
        emphasis_box = f'<div style="width:260px;flex-shrink:0;background:rgba(56,189,248,0.1);border-radius:16px;display:flex;align-items:center;justify-content:center;padding:20px;border:2px solid rgba(56,189,248,0.2);"><div style="text-align:center;"><div style="font-size:48px;font-weight:900;color:#38bdf8;">{emphasis[0]}</div><div style="font-size:16px;color:#94a3b8;margin-top:8px;">核心数据</div></div></div>'

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e293b 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative}}
</style></head><body>
<div style="position:absolute;top:0;left:0;right:0;height:80px;background:linear-gradient(90deg,#1e3a5f,#2563eb);display:flex;align-items:center;padding:0 80px;"><div style="font-size:34px;font-weight:700;color:white;">{title}</div></div>
<div style="position:absolute;top:100px;left:80px;right:80px;bottom:50px;display:flex;gap:30px;">
<div style="width:6px;background:linear-gradient(180deg,#38bdf8,#8b5cf6);border-radius:3px;flex-shrink:0;"></div>
<div style="flex:1;display:flex;flex-direction:column;justify-content:center;">{items_html}</div>
{emphasis_box}
</div>
</body></html>"""


def _render_comparison(slide: dict) -> str:
    """模板3: 对比 — 红绿双栏"""
    c = slide.get("content", {})
    title = c.get("title", "")
    comp = c.get("comparison", {})
    left_title = comp.get("left_title", "传统方案")
    left_items = comp.get("left_items", [])
    right_title = comp.get("right_title", "灵创方案")
    right_items = comp.get("right_items", [])
    emphasis = c.get("emphasis_words", [])

    left_html = "".join(
        f'<div style="display:flex;gap:10px;padding:10px 16px;font-size:22px;color:#fca5a5;line-height:1.5;"><span style="color:#ef4444;font-size:18px;margin-top:3px;">&#10007;</span>{_highlight(item, emphasis)}</div>'
        for item in left_items
    )
    right_html = "".join(
        f'<div style="display:flex;gap:10px;padding:10px 16px;font-size:22px;color:#86efac;line-height:1.5;"><span style="color:#22c55e;font-size:18px;margin-top:3px;">&#10003;</span>{_highlight(item, emphasis)}</div>'
        for item in right_items
    )

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e293b 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative}}
</style></head><body>
<div style="position:absolute;top:0;left:0;right:0;height:80px;background:linear-gradient(90deg,#1e3a5f,#2563eb);display:flex;align-items:center;padding:0 80px;"><div style="font-size:34px;font-weight:700;color:white;">{title}</div></div>
<div style="position:absolute;top:100px;left:80px;right:80px;bottom:50px;display:flex;gap:40px;">
<div style="flex:1;background:rgba(239,68,68,0.08);border-radius:16px;padding:24px;border:1px solid rgba(239,68,68,0.15);">
<div style="background:#ef4444;border-radius:10px;padding:12px 20px;margin-bottom:16px;"><div style="font-size:22px;font-weight:700;color:white;">{left_title}</div></div>
{left_html}
</div>
<div style="flex:1;background:rgba(34,197,94,0.08);border-radius:16px;padding:24px;border:1px solid rgba(34,197,94,0.15);">
<div style="background:#22c55e;border-radius:10px;padding:12px 20px;margin-bottom:16px;"><div style="font-size:22px;font-weight:700;color:white;">{right_title}</div></div>
{right_html}
</div>
</div>
</body></html>"""


def _render_big_number(slide: dict) -> str:
    """模板4: 大数字 — 超大渐变数字+说明"""
    c = slide.get("content", {})
    title = c.get("title", "")
    items = c.get("body_items", [])
    emphasis = c.get("emphasis_words", [])

    main_number = emphasis[0] if emphasis else (items[0][:20] if items else "100%")
    desc = items[0] if items else ""

    # 额外数据点
    extra_html = ""
    for item in items[1:4]:
        highlighted = _highlight(item, emphasis)
        extra_html += f'<div style="font-size:20px;color:#94a3b8;margin-top:12px;">{highlighted}</div>'

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e3a5f 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative;display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center}}
</style></head><body>
<div style="font-size:24px;color:#94a3b8;margin-bottom:20px;text-transform:uppercase;letter-spacing:4px;">{title}</div>
<div style="font-size:120px;font-weight:900;background:linear-gradient(135deg,#38bdf8,#a78bfa);-webkit-background-clip:text;-webkit-text-fill-color:transparent;line-height:1;">{main_number}</div>
<div style="font-size:22px;color:#94a3b8;margin-top:16px;">{desc}</div>
<div style="margin-top:30px;">{extra_html}</div>
</body></html>"""


def _render_quote(slide: dict) -> str:
    """模板5: 金句 — 深色背景+大号引用"""
    c = slide.get("content", {})
    title = c.get("title", "")
    items = c.get("body_items", [])
    emphasis = c.get("emphasis_words", [])

    quote = items[0] if items else title
    highlighted = _highlight(quote, emphasis)

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e3a5f 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative;display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;padding:0 160px}}
</style></head><body>
<div style="font-size:72px;color:#38bdf8;opacity:0.3;margin-bottom:-20px;">"</div>
<div style="font-size:40px;font-style:italic;line-height:1.6;color:#f1f5f9;">{highlighted}</div>
<div style="width:60px;height:3px;background:linear-gradient(90deg,#38bdf8,#8b5cf6);border-radius:2px;margin:30px 0 16px;"></div>
<div style="font-size:20px;color:#64748b;">{title}</div>
</body></html>"""


def _render_grid_3(slide: dict) -> str:
    """模板6: 三列卡片 — 三大优势/步骤"""
    c = slide.get("content", {})
    title = c.get("title", "")
    items = c.get("body_items", [])
    emphasis = c.get("emphasis_words", [])

    colors = ["#3b82f6", "#8b5cf6", "#06b6d4", "#f59e0b"]
    cards_html = ""
    for i, item in enumerate(items[:4]):
        color = colors[i % len(colors)]
        highlighted = _highlight(item, emphasis)
        cards_html += f"""<div style="flex:1;background:rgba(255,255,255,0.04);border-radius:16px;padding:24px;border-top:5px solid {color};">
<div style="font-size:48px;font-weight:900;color:{color};margin-bottom:12px;">{i + 1}</div>
<div style="font-size:20px;color:#e2e8f0;line-height:1.5;">{highlighted}</div>
</div>"""

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e293b 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative}}
</style></head><body>
<div style="position:absolute;top:0;left:0;right:0;height:80px;background:linear-gradient(90deg,#1e3a5f,#2563eb);display:flex;align-items:center;padding:0 80px;"><div style="font-size:34px;font-weight:700;color:white;">{title}</div></div>
<div style="position:absolute;top:110px;left:80px;right:80px;bottom:50px;display:flex;gap:24px;align-items:stretch;">
{cards_html}
</div>
</body></html>"""


def _render_closing(slide: dict) -> str:
    """模板7: 致谢页"""
    c = slide.get("content", {})
    title = c.get("title", "感谢聆听")
    items = c.get("body_items", [])

    contacts_html = "".join(
        f'<div style="font-size:18px;color:#94a3b8;margin-top:8px;">{item}</div>'
        for item in items
    )

    return f"""<!DOCTYPE html><html><head><meta charset="utf-8"><style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{width:{W}px;height:{H}px;background:linear-gradient(135deg,#0f172a 0%,#1e3a5f 100%);font-family:'PingFang SC','Microsoft YaHei',sans-serif;overflow:hidden;position:relative;display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center}}
</style></head><body>
<div style="position:absolute;width:500px;height:500px;border-radius:50%;background:radial-gradient(circle,rgba(56,189,248,0.06),transparent 70%);top:50%;left:50%;transform:translate(-50%,-50%);"></div>
<div style="font-size:64px;font-weight:900;background:linear-gradient(135deg,#38bdf8,#a78bfa);-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:20px;position:relative;z-index:1;">{title}</div>
<div style="width:80px;height:3px;background:linear-gradient(90deg,#38bdf8,#8b5cf6);border-radius:2px;margin-bottom:20px;position:relative;z-index:1;"></div>
<div style="font-size:22px;color:#94a3b8;margin-bottom:30px;position:relative;z-index:1;">期待与您合作</div>
<div style="position:relative;z-index:1;">{contacts_html}</div>
</body></html>"""


# ════════════════════════════════════════════════════════════════════
# 路由器
# ════════════════════════════════════════════════════════════════════

RENDERERS = {
    "cover": _render_cover,
    "bullet_points": _render_bullet_points,
    "comparison": _render_comparison,
    "big_number": _render_big_number,
    "quote": _render_quote,
    "grid_3": _render_grid_3,
    "grid_2": _render_bullet_points,
    "section_divider": _render_quote,
    "summary": _render_closing,
    "closing": _render_closing,
}


def _slide_to_html(slide: dict) -> str:
    layout = slide.get("layout_type", "bullet_points")
    renderer = RENDERERS.get(layout, _render_bullet_points)
    return renderer(slide)


# ════════════════════════════════════════════════════════════════════
# 渲染 + PPTX
# ════════════════════════════════════════════════════════════════════


def render_slides_to_images(slides: List[dict], output_dir: str) -> List[str]:
    """将 slides 渲染为 HTML 截图"""
    os.makedirs(output_dir, exist_ok=True)
    image_paths = []

    for i, slide in enumerate(slides):
        html = _slide_to_html(slide)
        html_path = os.path.join(output_dir, f"slide_{i:02d}.html")
        img_path = os.path.join(output_dir, f"slide_{i:02d}.png")

        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html)

        script = f"""
from playwright.sync_api import sync_playwright
with sync_playwright() as p:
    b = p.chromium.launch(headless=True, args=['--no-sandbox'])
    page = b.new_page(viewport={{'width': {W}, 'height': {H}}})
    page.goto('file:///{html_path.replace(chr(92), "/")}', wait_until='networkidle', timeout=10000)
    page.screenshot(path='{img_path.replace(chr(92), "/")}', full_page=False)
    b.close()
print('OK')
"""
        result = subprocess.run(
            ["python", "-c", script], capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0 and os.path.exists(img_path):
            image_paths.append(img_path)
            logger.info(f"Slide {i + 1}: {slide.get('layout_type', '?')}")
        else:
            logger.error(f"Slide {i + 1} screenshot failed: {result.stderr[:200]}")

    return image_paths


def images_to_pptx(
    image_paths: List[str], output_path: str, title: str = "Presentation"
):
    """将截图嵌入 PPTX"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title

    for img_path in image_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            img_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )

    prs.save(output_path)
