"""
PPT Parser - Extract structure, content, and design from reference PPTX
Based on ppt-master workflow for reference reconstruction
"""

import os
import json
from pathlib import Path
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import Counter


class PPTParser:
    """Parse reference PPTX to extract design spec and content outline"""

    def __init__(self, pptx_path: str):
        self.pptx_path = Path(pptx_path)
        self.prs = Presentation(str(self.pptx_path))
        self.design_spec = {}
        self.content_outline = []

    def parse(self) -> Dict[str, Any]:
        """Main parsing entry point"""
        print(f"Parsing reference PPT: {self.pptx_path}")

        # Extract canvas format
        canvas = self._extract_canvas_format()

        # Extract color scheme
        colors = self._extract_color_scheme()

        # Extract typography
        typography = self._extract_typography()

        # Extract content outline
        content = self._extract_content_outline()

        # Extract images
        images = self._extract_images()

        result = {
            "canvas": canvas,
            "colors": colors,
            "typography": typography,
            "content_outline": content,
            "images": images,
            "metadata": {
                "source_file": str(self.pptx_path),
                "total_slides": len(self.prs.slides),
                "slide_width": self.prs.slide_width,
                "slide_height": self.prs.slide_height,
            },
        }

        print(f"[OK] Parsed {len(self.prs.slides)} slides")
        return result

    def _extract_canvas_format(self) -> Dict[str, Any]:
        """Extract canvas dimensions and format"""
        width = self.prs.slide_width
        height = self.prs.slide_height
        ratio = width / height

        # Determine format
        if abs(ratio - 16 / 9) < 0.01:
            format_name = "ppt169"
        elif abs(ratio - 4 / 3) < 0.01:
            format_name = "ppt43"
        else:
            format_name = "custom"

        return {
            "format": format_name,
            "width_px": int(width / 9525),  # EMU to pixels
            "height_px": int(height / 9525),
            "width_inches": width / 914400,
            "height_inches": height / 914400,
            "aspect_ratio": ratio,
        }

    def _extract_color_scheme(self) -> Dict[str, Any]:
        """Extract color palette from all slides"""
        colors = []

        for slide in self.prs.slides:
            # Extract background color
            if slide.background.fill.type == 1:  # Solid fill
                try:
                    rgb = slide.background.fill.fore_color.rgb
                    colors.append(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass

            # Extract shape colors
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "fill") and shape.fill.type == 1:
                        rgb = shape.fill.fore_color.rgb
                        colors.append(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass

                try:
                    if hasattr(shape, "line") and shape.line.fill.type == 1:
                        rgb = shape.line.color.rgb
                        colors.append(f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass

        # Count and sort colors
        color_counts = Counter(colors)
        sorted_colors = sorted(color_counts.items(), key=lambda x: x[1], reverse=True)

        # Get top colors
        primary_colors = [c[0] for c in sorted_colors[:6]]

        return {
            "primary_colors": primary_colors,
            "all_colors": list(color_counts.keys()),
            "color_distribution": dict(sorted_colors[:15]),
        }

    def _extract_typography(self) -> Dict[str, Any]:
        """Extract font families and sizes"""
        fonts = []
        font_sizes = []

        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.append(run.font.name)
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)

        font_counts = Counter(fonts)
        size_counts = Counter(font_sizes)

        return {
            "primary_font": font_counts.most_common(1)[0][0]
            if font_counts
            else "Arial",
            "all_fonts": list(font_counts.keys()),
            "font_distribution": dict(font_counts.most_common(10)),
            "common_sizes": sorted(list(set(font_sizes)), reverse=True)[:10]
            if font_sizes
            else [24, 18, 14],
        }

    def _extract_content_outline(self) -> List[Dict[str, Any]]:
        """Extract content structure from each slide"""
        outline = []

        for idx, slide in enumerate(self.prs.slides, 1):
            slide_data = {
                "page_number": idx,
                "title": "",
                "body_text": [],
                "elements": [],
                "layout_type": self._detect_layout_type(slide),
                "element_count": len(slide.shapes),
            }

            # Extract text content
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                text = shape.text.strip()
                if not text:
                    continue

                # Detect if title (usually larger font, top position)
                if shape.top < self.prs.slide_height * 0.2:
                    if not slide_data["title"]:
                        slide_data["title"] = text
                else:
                    slide_data["body_text"].append(text)

                # Record element info
                slide_data["elements"].append(
                    {
                        "type": "text",
                        "content": text[:100],
                        "position": {
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        },
                    }
                )

            # Extract non-text elements
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_data["elements"].append(
                        {
                            "type": "image",
                            "position": {
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height,
                            },
                        }
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    slide_data["elements"].append(
                        {
                            "type": "table",
                            "rows": len(shape.table.rows),
                            "cols": len(shape.table.columns),
                        }
                    )
                elif shape.shape_type in [
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.FREEFORM,
                ]:
                    slide_data["elements"].append(
                        {"type": "shape", "shape_type": str(shape.shape_type)}
                    )

            outline.append(slide_data)

        return outline

    def _detect_layout_type(self, slide) -> str:
        """Detect slide layout type based on content"""
        text_shapes = [s for s in slide.shapes if hasattr(s, "text_frame")]
        image_shapes = [
            s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]

        if len(text_shapes) == 1 and len(image_shapes) == 0:
            return "title_only"
        elif len(text_shapes) >= 1 and len(image_shapes) >= 1:
            return "text_image"
        elif len(text_shapes) >= 2:
            return "text_heavy"
        elif len(image_shapes) >= 1:
            return "image_heavy"
        else:
            return "mixed"

    def _extract_images(self) -> List[Dict[str, Any]]:
        """Extract image information"""
        images = []

        for idx, slide in enumerate(self.prs.slides, 1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(
                        {
                            "slide_number": idx,
                            "position": {
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height,
                            },
                        }
                    )

        return images

    def save_to_json(self, output_path: str):
        """Save parsed result to JSON"""
        result = self.parse()
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(result, f, indent=2, ensure_ascii=False)

        print(f"[OK] Saved parsed result to: {output_file}")
        return result


def parse_reference_ppt(
    pptx_path: str, output_json: Optional[str] = None
) -> Dict[str, Any]:
    """
    Parse reference PPT and extract design spec + content outline

    Args:
        pptx_path: Path to reference PPTX file
        output_json: Optional path to save JSON output

    Returns:
        Parsed structure dict
    """
    parser = PPTParser(pptx_path)

    if output_json:
        return parser.save_to_json(output_json)
    else:
        return parser.parse()


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python ppt_parser.py <pptx_path> [output_json]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_json = sys.argv[2] if len(sys.argv) > 2 else None

    result = parse_reference_ppt(pptx_path, output_json)

    if not output_json:
        print(json.dumps(result, indent=2, ensure_ascii=False))
