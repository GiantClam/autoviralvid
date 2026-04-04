"""Pure Python PPTX-to-PPTX structural/visual comparator.

Compares two PPTX files across five dimensions without requiring
PowerPoint COM or LibreOffice. Uses python-pptx for extraction.
"""

from __future__ import annotations

import io
import math
import re
from dataclasses import dataclass, field
from difflib import SequenceMatcher
from typing import Any, Dict, List, Optional, Set, Tuple


def _clamp(value: float, lo: float = 0.0, hi: float = 100.0) -> float:
    return max(lo, min(hi, float(value)))


def _normalize_text(value: str) -> str:
    if not value:
        return ""
    text = str(value).strip().lower()
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"[^\w\s\u4e00-\u9fff%+.,!?;:]", "", text)
    return text


def _rgb_to_tuple(hex_color: str) -> Tuple[int, int, int]:
    if not hex_color:
        return (0, 0, 0)
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    if len(hex_color) != 6:
        return (0, 0, 0)
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return (r, g, b)
    except (ValueError, IndexError):
        return (0, 0, 0)


def _tuple_to_rgb(rgb: Tuple[int, int, int]) -> str:
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def _color_distance(hex1: str, hex2: str) -> float:
    rgb1 = _rgb_to_tuple(hex1)
    rgb2 = _rgb_to_tuple(hex2)
    return math.sqrt(sum((a - b) ** 2 for a, b in zip(rgb1, rgb2)))


def _colors_close(hex1: str, hex2: str, threshold: float = 40.0) -> bool:
    return _color_distance(hex1, hex2) <= threshold


def _fuzzy_color_overlap(
    ref_colors: List[str], gen_colors: List[str], threshold: float = 30.0
) -> float:
    if not ref_colors and not gen_colors:
        return 1.0
    if not ref_colors or not gen_colors:
        return 0.0
    matched = 0
    for rc in ref_colors:
        for gc in gen_colors:
            if _colors_close(rc, gc, threshold=threshold):
                matched += 1
                break
    return matched / max(len(ref_colors), len(gen_colors))


def _jaccard_similarity(set1: Set[str], set2: Set[str]) -> float:
    if not set1 and not set2:
        return 1.0
    if not set1 or not set2:
        return 0.0
    intersection = len(set1 & set2)
    union = len(set1 | set2)
    return intersection / union if union > 0 else 0.0


def _lcs_length(seq1: List[str], seq2: List[str]) -> int:
    m, n = len(seq1), len(seq2)
    dp = [[0] * (n + 1) for _ in range(m + 1)]
    for i in range(1, m + 1):
        for j in range(1, n + 1):
            if seq1[i - 1] == seq2[j - 1]:
                dp[i][j] = dp[i - 1][j - 1] + 1
            else:
                dp[i][j] = max(dp[i - 1][j], dp[i][j - 1])
    return dp[m][n]


def _mean(values: List[float]) -> float:
    return sum(values) / len(values) if values else 0.0


def _median(values: List[float]) -> float:
    if not values:
        return 0.0
    sorted_vals = sorted(values)
    n = len(sorted_vals)
    mid = n // 2
    if n % 2 == 0:
        return (sorted_vals[mid - 1] + sorted_vals[mid]) / 2
    return sorted_vals[mid]


# ── Data Models ──────────────────────────────────────────────


@dataclass
class ElementFingerprint:
    element_type: str
    left_norm: float
    top_norm: float
    width_norm: float
    height_norm: float
    text_content: str = ""
    fill_color: Optional[str] = None
    font_family: Optional[str] = None
    font_size_pt: Optional[float] = None
    text_color: Optional[str] = None


@dataclass
class SlideFingerprint:
    page_number: int
    inferred_type: str
    title: str
    body_texts: List[str]
    element_count: int
    elements: List[ElementFingerprint]
    dominant_colors: List[str] = field(default_factory=list)
    font_families: List[str] = field(default_factory=list)
    font_sizes: List[float] = field(default_factory=list)
    background_color: Optional[str] = None
    image_count: int = 0


@dataclass
class DeckFingerprint:
    slide_count: int
    slide_width_emu: int
    slide_height_emu: int
    aspect_ratio: float
    slides: List[SlideFingerprint]
    global_colors: List[str] = field(default_factory=list)
    global_font_families: List[str] = field(default_factory=list)
    theme_colors: List[str] = field(default_factory=list)
    media_count: int = 0
    total_image_count: int = 0
    deck_title: str = ""


@dataclass
class SlideComparisonDetail:
    ref_page: int
    gen_page: int
    title_similarity: float
    body_similarity: float
    element_count_ratio: float
    position_score: float
    color_overlap: float
    slide_score: float


@dataclass
class ComparisonReport:
    overall_score: float
    structure_score: float
    content_score: float
    visual_style_score: float
    geometry_score: float
    metadata_score: float
    slide_details: List[SlideComparisonDetail]
    issues: List[str] = field(default_factory=list)
    diagnostics: Dict[str, Any] = field(default_factory=dict)


# ── Extraction Functions ──────────────────────────────────────


def _infer_slide_type(slide_data: Dict[str, Any], idx: int, total: int) -> str:
    """Infer slide type from content and position."""
    title = slide_data.get("title", "").lower()
    body = " ".join(slide_data.get("body_texts", [])).lower()

    if idx == 0:
        if any(kw in title for kw in ["目录", "content", "contents", "index"]):
            return "toc"
        return "cover"

    if idx == total - 1:
        return "summary"

    if any(
        kw in title
        for kw in ["part", "section", "divider", "第一", "第二", "第三", "第四"]
    ):
        return "divider"

    if any(kw in title for kw in ["目录", "content", "contents"]):
        return "toc"

    return "content"


def _extract_element(
    slide: Any, shape: Any, slide_width_emu: int, slide_height_emu: int
) -> Optional[ElementFingerprint]:
    """Extract element fingerprint from a shape."""
    try:
        if not hasattr(shape, "shape_type"):
            return None

        shape_type = shape.shape_type
        element_type = "shape"

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            element_type = "text"
        elif hasattr(shape, "has_table") and shape.has_table:
            element_type = "table"
        elif hasattr(shape, "has_chart") and shape.has_chart:
            element_type = "chart"
        elif shape_type == 13:
            element_type = "image"

        left = int(getattr(shape, "left", 0) or 0)
        top = int(getattr(shape, "top", 0) or 0)
        width = int(getattr(shape, "width", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)

        left_norm = left / slide_width_emu if slide_width_emu > 0 else 0.0
        top_norm = top / slide_height_emu if slide_height_emu > 0 else 0.0
        width_norm = width / slide_width_emu if slide_width_emu > 0 else 0.0
        height_norm = height / slide_height_emu if slide_height_emu > 0 else 0.0

        text_content = ""
        fill_color = None
        font_family = None
        font_size_pt = None
        text_color = None

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            try:
                text_content = shape.text_frame.text or ""
            except:
                pass

            try:
                if hasattr(shape, "fill") and shape.fill:
                    if hasattr(shape.fill, "fore_color") and shape.fill.fore_color:
                        if (
                            hasattr(shape.fill.fore_color, "rgb")
                            and shape.fill.fore_color.rgb
                        ):
                            fill_color = f"#{shape.fill.fore_color.rgb}"
            except:
                pass

            try:
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        run = para.runs[0]
                        if hasattr(run, "font") and run.font:
                            if hasattr(run.font, "name") and run.font.name:
                                font_family = run.font.name
                            if hasattr(run.font, "size") and run.font.size:
                                font_size_pt = run.font.size.pt
                            if hasattr(run.font, "color") and run.font.color:
                                if (
                                    hasattr(run.font.color, "rgb")
                                    and run.font.color.rgb
                                ):
                                    text_color = f"#{run.font.color.rgb}"
            except:
                pass

        return ElementFingerprint(
            element_type=element_type,
            left_norm=left_norm,
            top_norm=top_norm,
            width_norm=width_norm,
            height_norm=height_norm,
            text_content=text_content.strip(),
            fill_color=fill_color,
            font_family=font_family,
            font_size_pt=font_size_pt,
            text_color=text_color,
        )
    except Exception:
        return None


def _extract_slide_fingerprint(
    slide: Any,
    page_number: int,
    total: int,
    slide_width_emu: int,
    slide_height_emu: int,
) -> SlideFingerprint:
    """Extract fingerprint from a single slide."""
    elements = []
    texts = []
    colors = []
    font_families = []
    font_sizes = []
    image_count = 0

    try:
        for shape in slide.shapes:
            elem = _extract_element(slide, shape, slide_width_emu, slide_height_emu)
            if elem:
                elements.append(elem)
                if elem.element_type == "image":
                    image_count += 1
                if elem.text_content:
                    texts.append(elem.text_content)
                if elem.fill_color:
                    colors.append(elem.fill_color)
                if elem.text_color:
                    colors.append(elem.text_color)
                if elem.font_family:
                    font_families.append(elem.font_family)
                if elem.font_size_pt:
                    font_sizes.append(elem.font_size_pt)
    except Exception:
        pass

    title = texts[0] if texts else ""
    body_texts = texts[1:] if len(texts) > 1 else []

    color_counts = {}
    for c in colors:
        if c:
            color_counts[c] = color_counts.get(c, 0) + 1
    dominant_colors = sorted(
        color_counts.keys(), key=lambda x: color_counts[x], reverse=True
    )[:5]

    unique_fonts = sorted(set(font_families))
    unique_sizes = sorted(set(font_sizes))

    inferred_type = _infer_slide_type(
        {"title": title, "body_texts": body_texts}, page_number - 1, total
    )

    return SlideFingerprint(
        page_number=page_number,
        inferred_type=inferred_type,
        title=title,
        body_texts=body_texts,
        element_count=len(elements),
        elements=elements,
        dominant_colors=dominant_colors,
        font_families=unique_fonts,
        font_sizes=unique_sizes,
        background_color=None,
        image_count=image_count,
    )


def _extract_theme_colors(pptx_bytes: bytes) -> List[str]:
    """Extract theme colors from PPTX file by parsing all theme XML files."""
    import zipfile
    import xml.etree.ElementTree as ET

    colors = []
    system_color_map = {
        "window": "#FFFFFF",  # White background
        "windowText": "#000000",  # Black text
        "btnFace": "#F0F0F0",  # Button face
        "btnText": "#000000",  # Button text
        "highlight": "#0078D4",  # Highlight
        "highlightText": "#FFFFFF",
    }

    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zip_file:
            # Extract ALL theme files
            theme_files = [
                f for f in zip_file.namelist() if "theme" in f and f.endswith(".xml")
            ]

            for theme_file in theme_files:
                theme_xml = zip_file.read(theme_file)
                root = ET.fromstring(theme_xml)

                for elem in root.iter():
                    if "clrScheme" in elem.tag:
                        for clr in elem:
                            tag = clr.tag.split("}")[1] if "}" in clr.tag else clr.tag
                            for srgb in clr.iter():
                                if "srgbClr" in srgb.tag:
                                    val = srgb.get("val")
                                    if val:
                                        colors.append(f"#{val}")
                                    break
                                elif "sysClr" in srgb.tag:
                                    val = srgb.get("val")
                                    if val:
                                        mapped = system_color_map.get(
                                            val.lower(), "#808080"
                                        )
                                        colors.append(mapped)
                                    break
    except Exception:
        pass

    return colors


def _count_media_assets(pptx_bytes: bytes) -> int:
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as zip_file:
            return len(
                [
                    name
                    for name in zip_file.namelist()
                    if name.startswith("ppt/media/") and not name.endswith("/")
                ]
            )
    except Exception:
        return 0


def extract_deck_fingerprint(pptx_bytes: bytes) -> DeckFingerprint:
    """Extract a complete fingerprint from a PPTX file."""
    try:
        from pptx import Presentation
        import io

        prs = Presentation(io.BytesIO(pptx_bytes))

        slide_count = len(prs.slides)
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height
        aspect_ratio = (
            slide_width_emu / slide_height_emu if slide_height_emu > 0 else 1.0
        )

        slides = []
        all_colors = []
        all_fonts = []
        deck_title = ""
        total_image_count = 0

        for idx, slide in enumerate(prs.slides):
            slide_fp = _extract_slide_fingerprint(
                slide, idx + 1, slide_count, slide_width_emu, slide_height_emu
            )
            slides.append(slide_fp)
            all_colors.extend(slide_fp.dominant_colors)
            all_fonts.extend(slide_fp.font_families)
            total_image_count += int(slide_fp.image_count or 0)
            if idx == 0 and slide_fp.title:
                deck_title = slide_fp.title

        color_counts = {}
        for c in all_colors:
            color_counts[c] = color_counts.get(c, 0) + 1
        global_colors = sorted(
            color_counts.keys(), key=lambda x: color_counts[x], reverse=True
        )[:10]

        theme_colors = []
        for color in _extract_theme_colors(pptx_bytes):
            if color not in theme_colors:
                theme_colors.append(color)

        for tc in theme_colors:
            if tc not in global_colors:
                global_colors.append(tc)

        global_font_families = sorted(set(all_fonts))
        media_count = _count_media_assets(pptx_bytes)

        return DeckFingerprint(
            slide_count=slide_count,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            aspect_ratio=aspect_ratio,
            slides=slides,
            global_colors=global_colors[:15],  # Include more colors
            global_font_families=global_font_families,
            theme_colors=theme_colors[:15],
            media_count=media_count,
            total_image_count=total_image_count,
            deck_title=deck_title,
        )
    except Exception:
        return DeckFingerprint(
            slide_count=0,
            slide_width_emu=0,
            slide_height_emu=0,
            aspect_ratio=1.0,
            slides=[],
            global_colors=[],
            global_font_families=[],
            theme_colors=[],
            media_count=0,
            total_image_count=0,
            deck_title="",
        )


# ── Comparison Functions ──────────────────────────────────────


def _compare_structure(
    ref: DeckFingerprint, gen: DeckFingerprint
) -> Tuple[float, List[str]]:
    """Compare structure: slide count, type sequence, element counts."""
    issues = []

    slide_count_score = max(
        0, 100 - abs(ref.slide_count - gen.slide_count) / max(ref.slide_count, 1) * 200
    )

    ref_types = [s.inferred_type for s in ref.slides]
    gen_types = [s.inferred_type for s in gen.slides]
    lcs_len = _lcs_length(ref_types, gen_types)
    type_lcs = lcs_len / max(len(ref_types), len(gen_types), 1)
    type_sequence_score = type_lcs * 100

    paired_slides = list(zip(ref.slides, gen.slides))
    elem_count_scores = []
    for r, g in paired_slides:
        score = max(
            0,
            100
            - abs(r.element_count - g.element_count) / max(r.element_count, 1) * 100,
        )
        elem_count_scores.append(score)
    elem_count_score = _mean(elem_count_scores) if elem_count_scores else 100.0

    structure_score = (
        0.40 * slide_count_score + 0.35 * type_sequence_score + 0.25 * elem_count_score
    )

    if abs(ref.slide_count - gen.slide_count) > 0:
        issues.append(f"Slide count mismatch: {ref.slide_count} vs {gen.slide_count}")

    return _clamp(structure_score), issues


def _compare_content(
    ref: DeckFingerprint, gen: DeckFingerprint
) -> Tuple[float, List[str]]:
    """Compare content: text similarity."""
    issues = []

    paired_slides = list(zip(ref.slides, gen.slides))
    title_sims = []
    body_sims = []

    for r, g in paired_slides:
        title_sim = SequenceMatcher(
            None, _normalize_text(r.title), _normalize_text(g.title)
        ).ratio()
        title_sims.append(title_sim)

        ref_all = set()
        for t in r.body_texts:
            if t:
                ref_all.update(_normalize_text(t).split())
        gen_all = set()
        for t in g.body_texts:
            if t:
                gen_all.update(_normalize_text(t).split())
        body_sim = _jaccard_similarity(ref_all, gen_all)
        body_sims.append(body_sim)

    title_score = _mean(title_sims) * 100 if title_sims else 0.0
    body_score = _mean(body_sims) * 100 if body_sims else 0.0
    unmatched_penalty = (
        abs(ref.slide_count - gen.slide_count) / max(ref.slide_count, 1) * 30
    )

    content_score = _clamp(0.45 * title_score + 0.55 * body_score - unmatched_penalty)

    if title_score < 50:
        issues.append(f"Low title similarity: {title_score:.1f}%")
    if body_score < 30:
        issues.append(f"Low body text similarity: {body_score:.1f}%")

    return content_score, issues


def _compare_visual_style(
    ref: DeckFingerprint, gen: DeckFingerprint
) -> Tuple[float, List[str]]:
    """Compare visual style: colors, fonts, backgrounds."""
    issues = []
    color_overlap = _fuzzy_color_overlap(ref.global_colors, gen.global_colors, threshold=30.0)
    theme_overlap = _fuzzy_color_overlap(ref.theme_colors, gen.theme_colors, threshold=22.0)
    font_overlap = _jaccard_similarity(
        set(ref.global_font_families), set(gen.global_font_families)
    )

    ref_sizes = [s.font_sizes for s in ref.slides if s.font_sizes]
    gen_sizes = [s.font_sizes for s in gen.slides if s.font_sizes]
    ref_all_sizes = set()
    for sizes in ref_sizes:
        ref_all_sizes.update(sizes)
    gen_all_sizes = set()
    for sizes in gen_sizes:
        gen_all_sizes.update(sizes)
    font_size_sim = _jaccard_similarity(ref_all_sizes, gen_all_sizes)

    bg_match_count = 0
    paired_slides = list(zip(ref.slides, gen.slides))
    for r, g in paired_slides:
        if r.background_color and g.background_color:
            if _colors_close(r.background_color, g.background_color):
                bg_match_count += 1
        elif not r.background_color and not g.background_color:
            bg_match_count += 1
    bg_match_ratio = bg_match_count / max(len(paired_slides), 1)

    visual_style_score = (
        0.25 * color_overlap * 100
        + 0.30 * theme_overlap * 100
        + 0.20 * font_overlap * 100
        + 0.10 * font_size_sim * 100
        + 0.15 * bg_match_ratio * 100
    )

    if color_overlap < 0.45:
        issues.append(f"Low color overlap: {color_overlap * 100:.1f}%")
    if theme_overlap < 0.70:
        issues.append(f"Low theme color overlap: {theme_overlap * 100:.1f}%")
    if font_overlap < 0.5:
        issues.append(f"Low font overlap: {font_overlap * 100:.1f}%")

    return _clamp(visual_style_score), issues


def _compare_geometry(
    ref: DeckFingerprint, gen: DeckFingerprint
) -> Tuple[float, List[str]]:
    """Compare element geometry: position and size."""
    issues = []

    def match_elements(
        ref_elems: List[ElementFingerprint], gen_elems: List[ElementFingerprint]
    ) -> float:
        if not ref_elems and not gen_elems:
            return 100.0
        if not ref_elems or not gen_elems:
            return 0.0

        ref_centers = [
            (e.left_norm + e.width_norm / 2, e.top_norm + e.height_norm / 2)
            for e in ref_elems
        ]
        gen_centers = [
            (e.left_norm + e.width_norm / 2, e.top_norm + e.height_norm / 2)
            for e in gen_elems
        ]

        matched_pairs = []
        used_gen = set()

        for i, (rx, ry) in enumerate(ref_centers):
            best_j = -1
            best_dist = float("inf")
            for j, (gx, gy) in enumerate(gen_centers):
                if j in used_gen:
                    continue
                dist = math.sqrt((rx - gx) ** 2 + (ry - gy) ** 2)
                if dist < best_dist:
                    best_dist = dist
                    best_j = j
            if best_j >= 0 and best_dist < 0.5:
                matched_pairs.append((i, best_j, best_dist))
                used_gen.add(best_j)

        if not matched_pairs:
            return 0.0

        elem_scores = []
        for i, j, dist in matched_pairs:
            size_diff = (
                abs(ref_elems[i].width_norm - gen_elems[j].width_norm)
                + abs(ref_elems[i].height_norm - gen_elems[j].height_norm)
            ) / 2
            elem_score = max(0, 100 - dist * 200 - size_diff * 200)
            elem_scores.append(elem_score)

        unmatched_penalty = (
            (len(ref_elems) + len(gen_elems) - 2 * len(matched_pairs))
            / max(len(ref_elems), len(gen_elems), 1)
            * 50
        )
        slide_geom = _mean(elem_scores) - unmatched_penalty if elem_scores else 0.0

        return _clamp(slide_geom)

    paired_slides = list(zip(ref.slides, gen.slides))
    slide_geom_scores = []
    for r, g in paired_slides:
        score = match_elements(r.elements, g.elements)
        slide_geom_scores.append(score)

    geometry_score = _mean(slide_geom_scores) if slide_geom_scores else 0.0

    if geometry_score < 50:
        issues.append(f"Low geometry match: {geometry_score:.1f}%")

    return _clamp(geometry_score), issues


def _compare_metadata(
    ref: DeckFingerprint, gen: DeckFingerprint
) -> Tuple[float, List[str]]:
    """Compare metadata: dimensions, aspect ratio."""
    issues = []

    ar_diff = abs(ref.aspect_ratio - gen.aspect_ratio) / max(ref.aspect_ratio, 0.01)
    dimension_score = max(0, 100 - ar_diff * 300)

    count_ratio = min(ref.slide_count, gen.slide_count) / max(
        ref.slide_count, gen.slide_count, 1
    )
    count_score = count_ratio * 100

    media_ratio = min(ref.media_count, gen.media_count) / max(ref.media_count, gen.media_count, 1)
    image_ratio = min(ref.total_image_count, gen.total_image_count) / max(
        ref.total_image_count, gen.total_image_count, 1
    )
    asset_score = 0.5 * media_ratio * 100 + 0.5 * image_ratio * 100

    metadata_score = 0.45 * dimension_score + 0.35 * count_score + 0.20 * asset_score

    if abs(ref.aspect_ratio - gen.aspect_ratio) > 0.1:
        issues.append(
            f"Aspect ratio mismatch: {ref.aspect_ratio:.2f} vs {gen.aspect_ratio:.2f}"
        )
    if ref.slide_count != gen.slide_count:
        issues.append(f"Slide count mismatch: {ref.slide_count} vs {gen.slide_count}")
    if ref.total_image_count != gen.total_image_count:
        issues.append(
            f"Image count mismatch: {ref.total_image_count} vs {gen.total_image_count}"
        )
    if ref.media_count != gen.media_count:
        issues.append(f"Media asset mismatch: {ref.media_count} vs {gen.media_count}")

    return _clamp(metadata_score), issues


def compare_decks(
    ref: DeckFingerprint, gen: DeckFingerprint, max_slides: int = 20
) -> ComparisonReport:
    """Compare two deck fingerprints and produce a comparison report.

    Args:
        max_slides: Maximum number of slides to compare (from start of each deck)
    """
    # Limit slides to compare
    ref_limited = DeckFingerprint(
        slide_count=ref.slide_count,
        slide_width_emu=ref.slide_width_emu,
        slide_height_emu=ref.slide_height_emu,
        aspect_ratio=ref.aspect_ratio,
        slides=ref.slides[:max_slides],
        global_colors=ref.global_colors,
        global_font_families=ref.global_font_families,
        theme_colors=ref.theme_colors,
        media_count=ref.media_count,
        total_image_count=ref.total_image_count,
        deck_title=ref.deck_title,
    )
    gen_limited = DeckFingerprint(
        slide_count=gen.slide_count,
        slide_width_emu=gen.slide_width_emu,
        slide_height_emu=gen.slide_height_emu,
        aspect_ratio=gen.aspect_ratio,
        slides=gen.slides[:max_slides],
        global_colors=gen.global_colors,
        global_font_families=gen.global_font_families,
        theme_colors=gen.theme_colors,
        media_count=gen.media_count,
        total_image_count=gen.total_image_count,
        deck_title=gen.deck_title,
    )

    structure_score, structure_issues = _compare_structure(ref_limited, gen_limited)
    content_score, content_issues = _compare_content(ref_limited, gen_limited)
    visual_style_score, visual_issues = _compare_visual_style(ref_limited, gen_limited)
    geometry_score, geometry_issues = _compare_geometry(ref_limited, gen_limited)
    metadata_score, metadata_issues = _compare_metadata(ref_limited, gen_limited)

    raw_score = _clamp(
        0.12 * structure_score
        + 0.23 * content_score
        + 0.30 * visual_style_score
        + 0.15 * geometry_score
        + 0.20 * metadata_score
    )

    hard_penalty = 0.0
    penalty_notes: List[str] = []

    slide_gap = abs(ref.slide_count - gen.slide_count)
    if slide_gap > 0:
        slide_penalty = min(25.0, slide_gap * 4.0)
        hard_penalty += slide_penalty
        penalty_notes.append(
            f"Hard penalty(slide_count): -{slide_penalty:.1f} ({ref.slide_count} vs {gen.slide_count})"
        )

    if ref.total_image_count > 0:
        image_ratio = min(ref.total_image_count, gen.total_image_count) / max(
            ref.total_image_count, gen.total_image_count, 1
        )
        if image_ratio < 0.95:
            image_penalty = (0.95 - image_ratio) * 25.0
            hard_penalty += image_penalty
            penalty_notes.append(
                f"Hard penalty(image_count): -{image_penalty:.1f} ({ref.total_image_count} vs {gen.total_image_count})"
            )

    if ref.theme_colors:
        theme_overlap_global = _fuzzy_color_overlap(
            ref.theme_colors, gen.theme_colors, threshold=22.0
        )
        if theme_overlap_global < 0.80:
            theme_penalty = (0.80 - theme_overlap_global) * 30.0
            hard_penalty += theme_penalty
            penalty_notes.append(
                f"Hard penalty(theme_colors): -{theme_penalty:.1f} (overlap={theme_overlap_global*100:.1f}%)"
            )

    overall_score = _clamp(raw_score - hard_penalty)

    all_issues = (
        structure_issues
        + content_issues
        + visual_issues
        + geometry_issues
        + metadata_issues
        + penalty_notes
    )

    slide_details = []
    paired_slides = list(zip(ref_limited.slides, gen_limited.slides))
    for r, g in paired_slides:
        title_sim = SequenceMatcher(
            None, _normalize_text(r.title), _normalize_text(g.title)
        ).ratio()

        ref_all = set()
        for t in r.body_texts:
            if t:
                ref_all.update(_normalize_text(t).split())
        gen_all = set()
        for t in g.body_texts:
            if t:
                gen_all.update(_normalize_text(t).split())
        body_sim = _jaccard_similarity(ref_all, gen_all)

        elem_count_ratio = min(r.element_count, g.element_count) / max(
            r.element_count, g.element_count, 1
        )

        position_score = 1.0
        color_overlap = 1.0

        slide_score = _clamp(
            0.4 * title_sim * 100 + 0.3 * body_sim * 100 + 0.3 * elem_count_ratio * 100
        )

        slide_details.append(
            SlideComparisonDetail(
                ref_page=r.page_number,
                gen_page=g.page_number,
                title_similarity=title_sim,
                body_similarity=body_sim,
                element_count_ratio=elem_count_ratio,
                position_score=position_score,
                color_overlap=color_overlap,
                slide_score=slide_score,
            )
        )

    diagnostics = {
        "ref_slide_count": ref.slide_count,
        "gen_slide_count": gen.slide_count,
        "ref_aspect_ratio": ref.aspect_ratio,
        "gen_aspect_ratio": gen.aspect_ratio,
        "ref_color_count": len(ref.global_colors),
        "gen_color_count": len(gen.global_colors),
        "ref_font_count": len(ref.global_font_families),
        "gen_font_count": len(gen.global_font_families),
        "ref_theme_color_count": len(ref.theme_colors),
        "gen_theme_color_count": len(gen.theme_colors),
        "ref_media_count": ref.media_count,
        "gen_media_count": gen.media_count,
        "ref_image_count": ref.total_image_count,
        "gen_image_count": gen.total_image_count,
        "raw_score": raw_score,
        "hard_penalty": hard_penalty,
    }

    return ComparisonReport(
        overall_score=overall_score,
        structure_score=structure_score,
        content_score=content_score,
        visual_style_score=visual_style_score,
        geometry_score=geometry_score,
        metadata_score=metadata_score,
        slide_details=slide_details,
        issues=all_issues,
        diagnostics=diagnostics,
    )


def compare_pptx_files(ref_path: str, gen_path: str) -> ComparisonReport:
    """Compare two PPTX files and return a comparison report."""
    from pathlib import Path

    ref_bytes = Path(ref_path).read_bytes()
    gen_bytes = Path(gen_path).read_bytes()

    ref_fp = extract_deck_fingerprint(ref_bytes)
    gen_fp = extract_deck_fingerprint(gen_bytes)

    return compare_decks(ref_fp, gen_fp)
