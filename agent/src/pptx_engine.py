"""
PPTX 母版渲染引擎 — 用 python-pptx 填充预设模板

核心思想: 不让代码控制审美，让代码去填空
预设 6 种版式: cover, bullet_points, split_left_img, split_right_img, quote, comparison, big_number
"""

from __future__ import annotations

import io
import logging
import os
import base64
import zipfile
from datetime import datetime
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.schemas.slide_v5 import SlideContentV5

logger = logging.getLogger("pptx_template")


# ════════════════════════════════════════════════════════════════════
# 配色方案
# ════════════════════════════════════════════════════════════════════

THEMES = {
    "professional": {
        "primary": RGBColor(0x1E, 0x3A, 0x5F),
        "secondary": RGBColor(0x25, 0x63, 0xEB),
        "accent": RGBColor(0x38, 0xBD, 0xF8),
        "text": RGBColor(0x1E, 0x29, 0x3B),
        "text_light": RGBColor(0x64, 0x74, 0x8B),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_light": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_dark": RGBColor(0x0F, 0x17, 0x2A),
    },
    "dark": {
        "primary": RGBColor(0x38, 0xBD, 0xF8),
        "secondary": RGBColor(0x8B, 0x5C, 0xF6),
        "accent": RGBColor(0x22, 0xC5, 0x5E),
        "text": RGBColor(0xE2, 0xE8, 0xF0),
        "text_light": RGBColor(0x94, 0xA3, 0xB8),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "bg_light": RGBColor(0x1E, 0x29, 0x3B),
        "bg_dark": RGBColor(0x0F, 0x17, 0x2A),
    },
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_text(
    tf,
    text: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    alignment=PP_ALIGN.LEFT,
    font_name: str = "Microsoft YaHei",
):
    """设置文本框内容"""
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def _set_rich_text(
    tf,
    items: List[str],
    font_size: int,
    color: RGBColor,
    font_name: str = "Microsoft YaHei",
):
    """设置带要点的富文本 (支持 **bold** 标记)"""
    tf.clear()
    for item in items:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)

        # 处理 **bold** 标记
        parts = re.split(r"\*\*(.+?)\*\*", item)
        for j, part in enumerate(parts):
            run = p.add_run()
            run.text = part
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = font_name
            if j % 2 == 1:  # 奇数索引是 bold 内容
                run.font.bold = True

    # 删除第一个空段落
    if tf.paragraphs and not tf.paragraphs[0].text:
        tf._txBody.remove(tf.paragraphs[0]._p)


import re


_MD_HEADING_RE = re.compile(r"^\s{0,3}#{1,6}\s+(.+?)\s*$")
_MD_BULLET_RE = re.compile(r"^\s*(?:[-*+]|(?:\d+[.)]))\s+(.+?)\s*$")


def _add_shape_bg(slide, color: RGBColor, left: int, top: int, width: int, height: int):
    """添加背景色块"""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide, color: RGBColor, left: int, top: int, height: int = 300):
    """添加强调色竖条"""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════════════
# 版式渲染器
# ════════════════════════════════════════════════════════════════════

import httpx


def _download_image(url: str) -> Optional[bytes]:
    """下载图片"""
    try:
        r = httpx.get(url, timeout=30, follow_redirects=True)
        r.raise_for_status()
        return r.content
    except Exception:
        return None


def render_cover(prs: Presentation, slide_data: SlideContentV5, theme: dict):
    """封面版式"""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # 全幅深色背景
    _add_shape_bg(
        slide, theme["bg_dark"], Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )

    # 装饰圆
    from pptx.enum.shapes import MSO_SHAPE

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(12), Inches(-1.5), Inches(5), Inches(5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = theme["accent"]
    circle.fill.fore_color.brightness = 0.0
    circle.line.fill.background()

    # 主标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
    _set_text(txBox.text_frame, slide_data.content.title, 44, theme["white"], bold=True)

    # 装饰线
    _add_shape_bg(
        slide, theme["accent"], Inches(1), Inches(3.8), Inches(1), Inches(0.04)
    )

    # 副标题
    if slide_data.content.subtitle:
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(4.1), Inches(11), Inches(0.6)
        )
        _set_text(txBox.text_frame, slide_data.content.subtitle, 22, theme["accent"])


def render_bullet_points(prs: Presentation, slide_data: SlideContentV5, theme: dict):
    """要点列表版式"""
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)

    # 标题栏
    _add_shape_bg(
        slide, theme["primary"], Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )

    # 标题
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.12), Inches(11), Inches(0.65)
    )
    _set_text(txBox.text_frame, slide_data.content.title, 28, theme["white"], bold=True)

    # 强调色竖条
    _add_accent_bar(
        slide,
        theme["accent"],
        Inches(0.7),
        Inches(1.2),
        Inches(min(len(slide_data.content.body_text) * 0.55, 5.5)),
    )

    # 要点列表
    if slide_data.content.body_text:
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1.2), Inches(10.5), Inches(5.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_rich_text(tf, slide_data.content.body_text, 18, theme["text"])

    # 高亮数据框 (右上)
    if slide_data.emphasis_words:
        highlight = slide_data.emphasis_words[0]
        _add_shape_bg(
            slide, theme["accent"], Inches(10), Inches(1.2), Inches(3), Inches(1.2)
        )
        txBox = slide.shapes.add_textbox(
            Inches(10.1), Inches(1.35), Inches(2.8), Inches(0.9)
        )
        _set_text(
            txBox.text_frame,
            highlight,
            24,
            theme["primary"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )


def render_comparison(prs: Presentation, slide_data: SlideContentV5, theme: dict):
    """对比版式"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    comp = slide_data.content.comparison

    # 标题栏
    _add_shape_bg(
        slide, theme["primary"], Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.12), Inches(11), Inches(0.65)
    )
    _set_text(txBox.text_frame, slide_data.content.title, 28, theme["white"], bold=True)

    if not comp:
        return

    # 左栏标题
    _add_shape_bg(
        slide, theme["secondary"], Inches(0.7), Inches(1.2), Inches(5.5), Inches(0.5)
    )
    txBox = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.25), Inches(5.2), Inches(0.4)
    )
    _set_text(txBox.text_frame, comp.left_title, 18, theme["white"], bold=True)

    # 右栏标题
    _add_shape_bg(
        slide, theme["accent"], Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.5)
    )
    txBox = slide.shapes.add_textbox(
        Inches(7.35), Inches(1.25), Inches(5.2), Inches(0.4)
    )
    _set_text(txBox.text_frame, comp.right_title, 18, theme["white"], bold=True)

    # 左栏要点
    if comp.left_items:
        items = [f"\u2713 {item}" for item in comp.left_items]
        txBox = slide.shapes.add_textbox(
            Inches(0.85), Inches(1.9), Inches(5.2), Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_rich_text(tf, items, 16, theme["text"])

    # 右栏要点
    if comp.right_items:
        items = [f"\u2713 {item}" for item in comp.right_items]
        txBox = slide.shapes.add_textbox(
            Inches(7.35), Inches(1.9), Inches(5.2), Inches(4.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        _set_rich_text(tf, items, 16, theme["text"])


def render_quote(prs: Presentation, slide_data: SlideContentV5, theme: dict):
    """名言/金句版式"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # 深色背景
    _add_shape_bg(
        slide, theme["bg_dark"], Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )

    # 引号装饰
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(1), Inches(1))
    _set_text(txBox.text_frame, "\u201c", 72, theme["accent"])

    # 金句正文
    text = (
        slide_data.content.body_text[0]
        if slide_data.content.body_text
        else slide_data.content.title
    )
    # 去掉 ** 标记但保留内容
    clean_text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    _set_text(tf, clean_text, 32, theme["white"])

    # 高亮词汇
    if slide_data.emphasis_words:
        highlighted = clean_text
        for w in slide_data.emphasis_words:
            highlighted = highlighted.replace(w, f"【{w}】")
        _set_text(tf, highlighted, 32, theme["white"])

    # 装饰线
    _add_shape_bg(
        slide, theme["accent"], Inches(2), Inches(5.2), Inches(1.5), Inches(0.04)
    )

    # 出处
    if slide_data.content.subtitle:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(8), Inches(0.5))
        _set_text(
            txBox.text_frame,
            f"\u2014\u2014 {slide_data.content.subtitle}",
            16,
            theme["text_light"],
        )


def render_big_number(prs: Presentation, slide_data: SlideContentV5, theme: dict):
    """大数字版式"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # 标题栏
    _add_shape_bg(
        slide, theme["primary"], Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.12), Inches(11), Inches(0.65)
    )
    _set_text(txBox.text_frame, slide_data.content.title, 28, theme["white"], bold=True)

    bn = slide_data.content.big_number
    if bn:
        # 大数字
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
        _set_text(
            txBox.text_frame,
            bn.number,
            96,
            theme["secondary"],
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        # 单位
        if bn.unit:
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(11), Inches(0.6)
            )
            _set_text(
                txBox.text_frame,
                bn.unit,
                24,
                theme["text_light"],
                alignment=PP_ALIGN.CENTER,
            )

        # 说明
        if bn.description:
            txBox = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(1))
            _set_text(
                txBox.text_frame,
                bn.description,
                18,
                theme["text"],
                alignment=PP_ALIGN.CENTER,
            )


def render_split(
    prs: Presentation, slide_data: SlideContentV5, theme: dict, side: str = "left"
):
    """左右分栏版式 (图左文右 / 图右文左)"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # 标题栏
    _add_shape_bg(
        slide, theme["primary"], Inches(0), Inches(0), prs.slide_width, Inches(0.9)
    )
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.12), Inches(11), Inches(0.65)
    )
    _set_text(txBox.text_frame, slide_data.content.title, 28, theme["white"], bold=True)

    if side == "left":
        # 左图
        if slide_data.content.image_url:
            img_data = _download_image(slide_data.content.image_url)
            if img_data:
                slide.shapes.add_picture(
                    io.BytesIO(img_data),
                    Inches(0.7),
                    Inches(1.2),
                    Inches(5.5),
                    Inches(4.5),
                )
        # 右文
        if slide_data.content.body_text:
            txBox = slide.shapes.add_textbox(
                Inches(6.5), Inches(1.2), Inches(6), Inches(5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            _set_rich_text(tf, slide_data.content.body_text, 18, theme["text"])
    else:
        # 左文
        if slide_data.content.body_text:
            txBox = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.2), Inches(6), Inches(5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            _set_rich_text(tf, slide_data.content.body_text, 18, theme["text"])
        # 右图
        if slide_data.content.image_url:
            img_data = _download_image(slide_data.content.image_url)
            if img_data:
                slide.shapes.add_picture(
                    io.BytesIO(img_data),
                    Inches(6.5),
                    Inches(1.2),
                    Inches(5.5),
                    Inches(4.5),
                )


# ════════════════════════════════════════════════════════════════════
# 统一入口
# ════════════════════════════════════════════════════════════════════

RENDERERS = {
    "cover": render_cover,
    "bullet_points": render_bullet_points,
    "comparison": render_comparison,
    "quote": render_quote,
    "big_number": render_big_number,
    "split_left_img": lambda prs, s, t: render_split(prs, s, t, "left"),
    "split_right_img": lambda prs, s, t: render_split(prs, s, t, "right"),
}


def generate_pptx(
    slides: List[SlideContentV5],
    title: str,
    author: str,
    template_id: str = "professional",
) -> bytes:
    """
    使用母版模板生成 PPTX。

    Args:
        slides: 语义化幻灯片数据
        title: 演示文稿标题
        author: 作者
        template_id: 模板ID (professional / dark)

    Returns:
        PPTX 文件字节
    """
    theme = THEMES.get(template_id, THEMES["professional"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = author

    for slide_data in slides:
        renderer = RENDERERS.get(slide_data.layout_type, render_bullet_points)
        renderer(prs, slide_data, theme)

    # 输出为 bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _strip_html_text(text: str) -> str:
    cleaned = re.sub(r"<[^>]+>", " ", str(text or ""))
    cleaned = cleaned.replace("&nbsp;", " ").replace("&amp;", "&").replace("&lt;", "<").replace("&gt;", ">")
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip()


def _parse_markitdown_sections(markdown_text: str) -> List[Dict[str, Any]]:
    sections: List[Dict[str, Any]] = []
    current: Dict[str, Any] = {"title": "", "bullets": []}
    lines = str(markdown_text or "").splitlines()
    for raw in lines:
        line = str(raw or "").strip()
        if not line:
            continue
        heading = _MD_HEADING_RE.match(line)
        if heading:
            title = _strip_html_text(heading.group(1))
            if current.get("title") or current.get("bullets"):
                sections.append(current)
            current = {"title": title, "bullets": []}
            continue
        bullet = _MD_BULLET_RE.match(line)
        if bullet:
            item = _strip_html_text(bullet.group(1))
            if item:
                current.setdefault("bullets", []).append(item)
            continue
        plain = _strip_html_text(line)
        if plain:
            current.setdefault("bullets", []).append(plain)
    if current.get("title") or current.get("bullets"):
        sections.append(current)
    return sections


def _merge_markitdown_replacements(
    *,
    global_map: Dict[str, str],
    per_slide_maps: List[Dict[str, str]],
    markdown_text: str,
) -> tuple[Dict[str, str], List[Dict[str, str]]]:
    sections = _parse_markitdown_sections(markdown_text)
    if not sections:
        return global_map, per_slide_maps

    merged_global = dict(global_map)
    merged_per_slide = [dict(item if isinstance(item, dict) else {}) for item in per_slide_maps]
    for idx, local_map in enumerate(merged_per_slide):
        if idx >= len(sections):
            break
        section = sections[idx] if isinstance(sections[idx], dict) else {}
        section_title = _strip_html_text(str(section.get("title") or ""))
        bullets = [
            _strip_html_text(item)
            for item in (section.get("bullets") if isinstance(section.get("bullets"), list) else [])
            if _strip_html_text(item)
        ]
        local_title = str(local_map.get("title") or "").strip()
        auto_title = bool(re.fullmatch(r"slide\s+\d+", local_title, flags=re.IGNORECASE))
        if section_title and (not local_title or auto_title):
            local_map["title"] = section_title
        if bullets and not str(local_map.get("subtitle") or "").strip():
            local_map["subtitle"] = bullets[0]
        if bullets and not str(local_map.get("body") or "").strip():
            local_map["body"] = "\n".join(bullets[:6]).strip()
        for bullet_idx in range(1, 7):
            key = f"bullet_{bullet_idx}"
            if str(local_map.get(key) or "").strip():
                continue
            if bullet_idx - 1 < len(bullets):
                local_map[key] = bullets[bullet_idx - 1]
        slide_no = idx + 1
        for key in ("title", "subtitle", "body", "narration", "speaker_notes"):
            value = str(local_map.get(key) or "").strip()
            if value:
                merged_global[f"slide_{slide_no}_{key}"] = value
        for bullet_idx in range(1, 7):
            bullet_value = str(local_map.get(f"bullet_{bullet_idx}") or "").strip()
            if bullet_value:
                merged_global[f"slide_{slide_no}_bullet_{bullet_idx}"] = bullet_value
    return merged_global, merged_per_slide


def _extract_slide_image_url(slide: Dict[str, Any]) -> str:
    def _norm_url(value: Any) -> str:
        text = str(value or "").strip()
        if not text:
            return ""
        lowered = text.lower()
        if lowered.startswith(("http://", "https://", "data:image/")):
            return text
        return ""

    for key in ("image_url", "cover_image_url", "hero_image_url"):
        hit = _norm_url(slide.get(key))
        if hit:
            return hit

    for container_key in ("elements", "blocks"):
        container = slide.get(container_key)
        if not isinstance(container, list):
            continue
        for item in container:
            if not isinstance(item, dict):
                continue
            content = item.get("content")
            if isinstance(content, dict):
                for key in ("url", "image_url", "src"):
                    hit = _norm_url(content.get(key))
                    if hit:
                        return hit
            hit = _norm_url(item.get("url") or item.get("image_url"))
            if hit:
                return hit
            if isinstance(content, str):
                hit = _norm_url(content)
                if hit:
                    return hit

    return ""


def _replace_text_tokens(text: str, mapping: Dict[str, str]) -> tuple[str, int]:
    out = str(text or "")
    replaced = 0
    for key, value in mapping.items():
        if not key:
            continue
        pattern = re.compile(r"\{\{\s*" + re.escape(str(key)) + r"\s*\}\}", flags=re.IGNORECASE)
        out, count = pattern.subn(str(value or ""), out)
        replaced += int(count)
    return out, replaced


def _replace_in_text_frame(text_frame, mapping: Dict[str, str]) -> int:
    replaced = 0
    for paragraph in text_frame.paragraphs:
        if paragraph.runs:
            for run in paragraph.runs:
                next_text, count = _replace_text_tokens(run.text or "", mapping)
                if count > 0:
                    run.text = next_text
                    replaced += count
        else:
            next_text, count = _replace_text_tokens(paragraph.text or "", mapping)
            if count > 0:
                paragraph.text = next_text
                replaced += count
    return replaced


def _extract_slide_index_from_xml_part(part_name: str) -> Optional[int]:
    name = str(part_name or "").replace("\\", "/").strip().lower()
    match = re.search(r"^ppt/slides/slide(\d+)\.xml$", name)
    if match:
        try:
            return int(match.group(1))
        except Exception:
            return None
    match = re.search(r"^ppt/notesslides/notesslide(\d+)\.xml$", name)
    if match:
        try:
            return int(match.group(1))
        except Exception:
            return None
    return None


def _replace_tokens_in_xml_bytes(xml_bytes: bytes, mapping: Dict[str, str]) -> tuple[bytes, int]:
    if not xml_bytes:
        return xml_bytes, 0
    try:
        text = xml_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = xml_bytes.decode("utf-8", errors="replace")
    replaced_text, replaced_count = _replace_text_tokens(text, mapping)
    if replaced_count <= 0:
        return xml_bytes, 0
    return replaced_text.encode("utf-8"), int(replaced_count)


def _is_xml_part(lowered: str) -> bool:
    return lowered.endswith(".xml") and (
        lowered.startswith("ppt/") or lowered.startswith("docprops/") or lowered.startswith("_rels/")
    )


def _looks_like_image_ref(value: Any) -> bool:
    text = str(value or "").strip().lower()
    return text.startswith(("http://", "https://", "data:image/"))


def _resolve_image_bytes(image_ref: str) -> Optional[bytes]:
    text = str(image_ref or "").strip()
    if not text:
        return None
    lowered = text.lower()
    if lowered.startswith(("http://", "https://")):
        return _download_image(text)
    if lowered.startswith("data:image/") and "," in text:
        _, encoded = text.split(",", 1)
        try:
            return base64.b64decode(encoded, validate=False)
        except Exception:
            return None
    return None


def _collect_referenced_media_names(parts: Dict[str, bytes]) -> set[str]:
    names: set[str] = set()
    patterns = [
        re.compile(r"(?:\.\./)?media/([^\"'>\s]+)", flags=re.IGNORECASE),
        re.compile(r"/ppt/media/([^\"'>\s]+)", flags=re.IGNORECASE),
    ]
    for part_name, payload in parts.items():
        lowered = str(part_name or "").replace("\\", "/").lower()
        if not (lowered.endswith(".xml") or lowered.endswith(".rels")):
            continue
        try:
            text = payload.decode("utf-8")
        except UnicodeDecodeError:
            text = payload.decode("utf-8", errors="ignore")
        for pattern in patterns:
            for match in pattern.finditer(text):
                candidate = str(match.group(1) or "").strip()
                if candidate:
                    names.add(candidate.split("/")[-1])
    return names


def _replace_image_placeholders_in_shape(slide, shape, mapping: Dict[str, str]) -> int:
    if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
        return 0
    tf = shape.text_frame
    paragraph_text = " ".join(str(paragraph.text or "") for paragraph in tf.paragraphs)
    if not paragraph_text:
        return 0

    matched_key = ""
    matched_ref = ""
    for key, value in mapping.items():
        image_ref = str(value or "").strip()
        if not image_ref or not _looks_like_image_ref(image_ref):
            continue
        token_pattern = re.compile(r"\{\{\s*" + re.escape(str(key)) + r"\s*\}\}", flags=re.IGNORECASE)
        if token_pattern.search(paragraph_text):
            matched_key = str(key)
            matched_ref = image_ref
            break

    if not matched_key:
        return 0
    image_bytes = _resolve_image_bytes(matched_ref)
    if not image_bytes:
        return 0

    try:
        slide.shapes.add_picture(
            io.BytesIO(image_bytes),
            shape.left,
            shape.top,
            shape.width,
            shape.height,
        )
        shape._element.getparent().remove(shape._element)
        return 1
    except Exception:
        return 0


def _template_contains_image_placeholders(template_bytes: bytes) -> bool:
    if not template_bytes:
        return False
    token_pattern = re.compile(
        r"\{\{\s*(?:image(?:_url)?|slide_\d+_image(?:_url)?|image_\d+)\s*\}\}",
        flags=re.IGNORECASE,
    )
    try:
        with zipfile.ZipFile(io.BytesIO(template_bytes), mode="r") as zin:
            for info in zin.infolist():
                lowered = str(info.filename or "").replace("\\", "/").lower()
                if not lowered.endswith(".xml"):
                    continue
                payload = zin.read(info.filename)
                try:
                    text = payload.decode("utf-8")
                except UnicodeDecodeError:
                    text = payload.decode("utf-8", errors="ignore")
                if token_pattern.search(text):
                    return True
    except Exception:
        return False
    return False


def _has_image_refs(mapping: Dict[str, str]) -> bool:
    for value in mapping.values():
        if _looks_like_image_ref(value):
            return True
    return False


def _fill_template_pptx_with_xml(
    *,
    template_bytes: bytes,
    global_map: Dict[str, str],
    per_slide: List[Dict[str, str]],
) -> Dict[str, Any]:
    in_mem = io.BytesIO(template_bytes)
    out_mem = io.BytesIO()
    total_replaced = 0
    template_slide_count = 0
    cleaned_resource_count = 0
    parts: Dict[str, bytes] = {}
    infos: Dict[str, zipfile.ZipInfo] = {}

    with zipfile.ZipFile(in_mem, mode="r") as zin:
        for info in zin.infolist():
            part_name = str(info.filename or "")
            lowered = part_name.replace("\\", "/").lower()
            if lowered.startswith("ppt/slides/slide"):
                template_slide_count += 1
            payload = zin.read(info.filename)
            if _is_xml_part(lowered):
                slide_idx = _extract_slide_index_from_xml_part(lowered)
                merged_map = dict(global_map)
                if slide_idx is not None and slide_idx > 0 and slide_idx - 1 < len(per_slide):
                    merged_map.update(per_slide[slide_idx - 1])
                payload, replaced = _replace_tokens_in_xml_bytes(payload, merged_map)
                total_replaced += int(replaced)
            parts[part_name] = payload
            infos[part_name] = info

    referenced_media = _collect_referenced_media_names(parts)
    media_entries = [
        part_name
        for part_name in parts.keys()
        if str(part_name or "").replace("\\", "/").lower().startswith("ppt/media/")
    ]
    cleanup_enabled = bool(media_entries)

    with zipfile.ZipFile(out_mem, mode="w", compression=zipfile.ZIP_DEFLATED) as zout:
        for part_name, payload in parts.items():
            lowered = str(part_name or "").replace("\\", "/").lower()
            if cleanup_enabled and lowered.startswith("ppt/media/"):
                media_name = lowered.split("/")[-1]
                if media_name not in referenced_media:
                    cleaned_resource_count += 1
                    continue
            zout.writestr(infos[part_name], payload)

    return {
        "pptx_bytes": out_mem.getvalue(),
        "replacement_count": int(total_replaced),
        "token_keys": sorted(set(str(k) for k in global_map.keys())),
        "slides_used": int(len(per_slide)),
        "template_slide_count": int(template_slide_count),
        "engine": "xml",
        "cleaned_resource_count": int(cleaned_resource_count),
    }


def _fill_template_pptx_with_python_pptx(
    *,
    template_bytes: bytes,
    global_map: Dict[str, str],
    per_slide: List[Dict[str, str]],
) -> Dict[str, Any]:
    prs = Presentation(io.BytesIO(template_bytes))
    total_replaced = 0
    image_replaced = 0
    token_keys = sorted(set(str(k) for k in global_map.keys()))

    for slide_idx, slide in enumerate(prs.slides):
        local_map = per_slide[slide_idx] if slide_idx < len(per_slide) and isinstance(per_slide[slide_idx], dict) else {}
        merged_map = {**global_map, **local_map}
        for shape in list(slide.shapes):
            image_replaced += _replace_image_placeholders_in_shape(slide, shape, merged_map)
            if not shape._element.getparent():
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                total_replaced += _replace_in_text_frame(shape.text_frame, merged_map)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        total_replaced += _replace_in_text_frame(cell.text_frame, merged_map)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            total_replaced += _replace_in_text_frame(slide.notes_slide.notes_text_frame, merged_map)

    out = io.BytesIO()
    prs.save(out)
    return {
        "pptx_bytes": out.getvalue(),
        "replacement_count": int(total_replaced + image_replaced),
        "token_keys": token_keys,
        "slides_used": int(len(per_slide)),
        "template_slide_count": int(len(prs.slides)),
        "engine": "python-pptx",
        "image_replacement_count": int(image_replaced),
        "cleaned_resource_count": 0,
    }


def fill_template_pptx(
    *,
    template_bytes: bytes,
    slides: List[Dict[str, Any]],
    deck_title: str,
    author: str,
) -> Dict[str, Any]:
    """
    Fill placeholders in an uploaded PPTX template.

    Supported placeholders:
    - deck-level: {{deck_title}}, {{title}}, {{deck_author}}, {{author}}, {{date}}, {{year}}
    - per-slide (global): {{slide_1_title}}, {{slide_1_body}}, {{slide_1_bullet_1}} ...
    - per-current-slide: {{title}}, {{subtitle}}, {{body}}, {{bullet_1}} ... {{bullet_6}}
    """
    if not template_bytes:
        raise ValueError("template_bytes is empty")
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides is empty")

    markitdown_enabled = str(os.getenv("PPT_TEMPLATE_MARKITDOWN_ENABLED", "true")).strip().lower() not in {
        "0",
        "false",
        "no",
        "off",
    }
    markitdown_ok = False
    markitdown_issue = ""
    markitdown_text = ""
    if not markitdown_enabled:
        raise ValueError("markitdown_template_strict_mode_requires_PPT_TEMPLATE_MARKITDOWN_ENABLED=true")
    try:
        from src.ppt_visual_qa import extract_text_with_markitdown  # lazy import to avoid hard dependency loop

        extracted = extract_text_with_markitdown(
            template_bytes,
            timeout_sec=max(5, int(os.getenv("PPT_TEMPLATE_MARKITDOWN_TIMEOUT_SEC", "20"))),
        )
        markitdown_ok = bool(extracted.get("ok"))
        markitdown_text = str(extracted.get("text") or "")
        markitdown_issue = str(extracted.get("error") or "")
    except Exception as exc:
        markitdown_ok = False
        markitdown_issue = f"markitdown_template_extract_failed: {str(exc)[:160]}"

    if not markitdown_ok or not markitdown_text.strip():
        raise ValueError(markitdown_issue or "markitdown_template_extract_empty")

    now = datetime.utcnow()
    global_map: Dict[str, str] = {
        "deck_title": str(deck_title or "").strip(),
        "title": str(deck_title or "").strip(),
        "deck_author": str(author or "").strip(),
        "author": str(author or "").strip(),
        "date": now.strftime("%Y-%m-%d"),
        "year": now.strftime("%Y"),
    }
    per_slide: List[Dict[str, str]] = []
    normalized_slides = [dict(item) for item in slides if isinstance(item, dict)]
    for idx, slide in enumerate(normalized_slides):
        image_url = _extract_slide_image_url(slide)
        slide_no = idx + 1
        local_map = {
            "slide_index": str(slide_no),
            "image": image_url,
            "image_url": image_url,
        }
        per_slide.append(local_map)
        global_map[f"slide_{slide_no}_image"] = image_url
        global_map[f"slide_{slide_no}_image_url"] = image_url
        global_map[f"image_{slide_no}"] = image_url

    global_map, per_slide = _merge_markitdown_replacements(
        global_map=global_map,
        per_slide_maps=per_slide,
        markdown_text=markitdown_text,
    )
    if not any(str((item or {}).get("title") or "").strip() for item in per_slide):
        raise ValueError("markitdown_template_sections_missing")
    template_has_image_tokens = _template_contains_image_placeholders(template_bytes)
    image_ref_available = _has_image_refs(global_map) or any(
        _has_image_refs(local_map if isinstance(local_map, dict) else {}) for local_map in per_slide
    )
    needs_image_fill = bool(template_has_image_tokens and image_ref_available)

    engine = str(os.getenv("PPT_TEMPLATE_EDIT_ENGINE", "xml")).strip().lower()
    if engine not in {"xml", "python-pptx", "auto"}:
        engine = "xml"
    if needs_image_fill and engine in {"xml", "auto"}:
        result = _fill_template_pptx_with_python_pptx(
            template_bytes=template_bytes,
            global_map=global_map,
            per_slide=per_slide,
        )
        result["markitdown_enabled"] = markitdown_enabled
        result["markitdown_ok"] = markitdown_ok
        result["markitdown_issue"] = markitdown_issue
        result["markitdown_used"] = bool(markitdown_ok and markitdown_text.strip())
        return result

    if engine in {"xml", "auto"}:
        try:
            xml_result = _fill_template_pptx_with_xml(
                template_bytes=template_bytes,
                global_map=global_map,
                per_slide=per_slide,
            )
            if int(xml_result.get("replacement_count") or 0) > 0 or engine == "xml":
                xml_result["markitdown_enabled"] = markitdown_enabled
                xml_result["markitdown_ok"] = markitdown_ok
                xml_result["markitdown_issue"] = markitdown_issue
                xml_result["markitdown_used"] = bool(markitdown_ok and markitdown_text.strip())
                return xml_result
        except Exception as exc:
            logger.warning("[pptx_template] xml template edit fallback to python-pptx: %s", exc)
            if engine == "xml":
                # xml mode is explicitly requested; still try python-pptx as hard fallback.
                pass

    result = _fill_template_pptx_with_python_pptx(
        template_bytes=template_bytes,
        global_map=global_map,
        per_slide=per_slide,
    )
    result["markitdown_enabled"] = markitdown_enabled
    result["markitdown_ok"] = markitdown_ok
    result["markitdown_issue"] = markitdown_issue
    result["markitdown_used"] = bool(markitdown_ok and markitdown_text.strip())
    return result
