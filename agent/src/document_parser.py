"""文档解析器 — Feature B: 解析PPT/PDF文件为SlideContent[]"""

from __future__ import annotations

import ipaddress
import logging
import os
import tempfile
from typing import Optional
from urllib.parse import urlparse

import httpx

from src.schemas.ppt import ParsedDocument, SlideBackground, SlideContent, SlideElement

logger = logging.getLogger("document_parser")

# 最大下载文件大小: 50MB
MAX_FILE_SIZE = 50 * 1024 * 1024

# 禁止的内网IP段 (SSRF防护)
_BLOCKED_NETWORKS = [
    ipaddress.ip_network("10.0.0.0/8"),
    ipaddress.ip_network("172.16.0.0/12"),
    ipaddress.ip_network("192.168.0.0/16"),
    ipaddress.ip_network("127.0.0.0/8"),
    ipaddress.ip_network("169.254.0.0/16"),
    ipaddress.ip_network("0.0.0.0/8"),
    ipaddress.ip_network("::1/128"),
    ipaddress.ip_network("fc00::/7"),
]


def _validate_url_safety(url: str) -> None:
    """SSRF防护: 校验URL不指向内网地址"""
    parsed = urlparse(url)
    hostname = parsed.hostname
    if not hostname:
        raise ValueError(f"无效URL: 无法解析主机名")

    try:
        import socket

        resolved = socket.getaddrinfo(hostname, None)
        for family, _, _, _, sockaddr in resolved:
            ip = ipaddress.ip_address(sockaddr[0])
            for blocked in _BLOCKED_NETWORKS:
                if ip in blocked:
                    raise ValueError(f"禁止访问内网地址: {hostname} ({ip})")
    except socket.gaierror:
        raise ValueError(f"无法解析主机名: {hostname}")


async def _download_file(url: str, suffix: str) -> str:
    """下载文件到临时目录，返回本地路径。带SSRF防护和大小限制。"""
    # SSRF防护
    _validate_url_safety(url)

    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    tmp.close()

    async with httpx.AsyncClient(timeout=120, follow_redirects=True) as client:
        # HEAD请求检查Content-Length
        try:
            head = await client.head(url)
            content_length = head.headers.get("content-length")
            if content_length and int(content_length) > MAX_FILE_SIZE:
                os.unlink(tmp.name)
                raise ValueError(
                    f"文件过大: {int(content_length) / 1024 / 1024:.1f}MB, "
                    f"最大允许 {MAX_FILE_SIZE / 1024 / 1024:.0f}MB"
                )
        except httpx.HTTPError:
            pass  # HEAD不支持时继续

        # 流式下载
        total_bytes = 0
        with open(tmp.name, "wb") as f:
            async with client.stream("GET", url) as resp:
                resp.raise_for_status()
                async for chunk in resp.aiter_bytes(chunk_size=8192):
                    total_bytes += len(chunk)
                    if total_bytes > MAX_FILE_SIZE:
                        f.close()
                        os.unlink(tmp.name)
                        raise ValueError(
                            f"文件过大: 超过 {MAX_FILE_SIZE / 1024 / 1024:.0f}MB 限制"
                        )
                    f.write(chunk)

    return tmp.name


def _extract_slide_title(slide) -> str:
    """从pptx slide中提取标题"""
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()[:100]
    return ""


def _sanitize_text(text: str) -> str:
    """HTML转义，防止XSS"""
    import html

    return html.escape(text)


def _shape_to_element(shape, index: int) -> Optional[SlideElement]:
    """将pptx shape转换为SlideElement"""
    from pptx.util import Emu

    left = int(Emu(shape.left).inches * 96) if hasattr(shape, "left") else 0
    top = int(Emu(shape.top).inches * 96) if hasattr(shape, "top") else 0
    width = int(Emu(shape.width).inches * 96) if hasattr(shape, "width") else 200
    height = int(Emu(shape.height).inches * 96) if hasattr(shape, "height") else 100

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if not text:
            return None
        return SlideElement(
            type="text",
            left=left,
            top=top,
            width=width,
            height=height,
            content=_sanitize_text(text),
            style={"fontSize": 18, "fontFamily": "Microsoft YaHei", "color": "#333333"},
        )

    if hasattr(shape, "image"):
        return SlideElement(
            type="image",
            left=left,
            top=top,
            width=width,
            height=height,
            src="",
            style={"objectFit": "cover"},
        )

    if shape.has_table:
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [_sanitize_text(cell.text) for cell in row.cells]
            rows.append(cells)
        return SlideElement(
            type="table",
            left=left,
            top=top,
            width=width,
            height=height,
            table_rows=rows,
            style={"fontSize": 14},
        )

    return None


def _parse_pptx_sync(file_path: str) -> ParsedDocument:
    """同步解析PPTX文件"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides: list[SlideContent] = []

    for i, slide in enumerate(prs.slides):
        title = _extract_slide_title(slide)
        elements: list[SlideElement] = []

        for j, shape in enumerate(slide.shapes):
            el = _shape_to_element(shape, j)
            if el:
                elements.append(el)

        if not elements and title:
            elements.append(
                SlideElement(
                    type="text",
                    left=100,
                    top=100,
                    width=1720,
                    height=80,
                    content=f"<b>{_sanitize_text(title)}</b>",
                    style={
                        "fontSize": 40,
                        "fontFamily": "Microsoft YaHei",
                        "color": "#333333",
                        "bold": True,
                    },
                )
            )

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(
            SlideContent(
                outline_id=f"parsed-{i}",
                order=i,
                title=title or f"第 {i + 1} 页",
                elements=elements,
                background=SlideBackground(type="solid", color="#ffffff"),
                narration=_sanitize_text(notes),
                speaker_notes=_sanitize_text(notes),
                duration=120,
            )
        )

    title = prs.core_properties.title or ""

    return ParsedDocument(
        source_type="pptx",
        source_url=file_path,
        title=title or "导入的PPT",
        slides=slides,
        total_pages=len(slides),
    )


def _parse_pdf_sync(file_path: str) -> ParsedDocument:
    """同步解析PDF文件"""
    import pdfplumber

    slides: list[SlideContent] = []

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            lines = [l.strip() for l in text.split("\n") if l.strip()]

            title = _sanitize_text(lines[0][:100]) if lines else f"第 {i + 1} 页"

            elements: list[SlideElement] = []

            if title:
                elements.append(
                    SlideElement(
                        type="text",
                        left=100,
                        top=60,
                        width=1720,
                        height=60,
                        content=f"<b>{title}</b>",
                        style={
                            "fontSize": 36,
                            "fontFamily": "Microsoft YaHei",
                            "color": "#333333",
                            "bold": True,
                        },
                    )
                )

            body = "\n".join(lines[1:]) if len(lines) > 1 else text
            if body.strip():
                elements.append(
                    SlideElement(
                        type="text",
                        left=100,
                        top=160,
                        width=1720,
                        height=700,
                        content=_sanitize_text(body),
                        style={
                            "fontSize": 20,
                            "fontFamily": "Microsoft YaHei",
                            "color": "#555555",
                        },
                    )
                )

            slides.append(
                SlideContent(
                    outline_id=f"parsed-{i}",
                    order=i,
                    title=title,
                    elements=elements,
                    background=SlideBackground(type="solid", color="#ffffff"),
                    narration=_sanitize_text(body[:500]) if body else "",
                    duration=120,
                )
            )

    return ParsedDocument(
        source_type="pdf",
        source_url=file_path,
        title=os.path.basename(file_path),
        slides=slides,
        total_pages=len(slides),
    )


async def parse_document(file_url: str, file_type: str = "pptx") -> ParsedDocument:
    """
    解析PPT/PDF文件为结构化数据。

    Args:
        file_url: 文件URL (R2 CDN 或 HTTP)
        file_type: "pptx" | "ppt" | "pdf"

    Returns:
        ParsedDocument

    Raises:
        ValueError: URL不安全或文件过大
    """
    import asyncio

    suffix_map = {"pptx": ".pptx", "ppt": ".ppt", "pdf": ".pdf"}
    suffix = suffix_map.get(file_type, ".pptx")

    logger.info(f"[document_parser] Downloading {file_type} from {file_url[:100]}...")
    local_path = await _download_file(file_url, suffix)

    try:
        if file_type == "pdf":
            return await asyncio.to_thread(_parse_pdf_sync, local_path)
        else:
            return await asyncio.to_thread(_parse_pptx_sync, local_path)
    finally:
        try:
            os.unlink(local_path)
        except Exception:
            pass
